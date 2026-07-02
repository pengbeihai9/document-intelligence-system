import hashlib
import io
import json
import os
import re
import shutil
import secrets
import sqlite3
import subprocess
import time
import unicodedata
import uuid
from collections import Counter
from pathlib import Path
from tempfile import NamedTemporaryFile, TemporaryDirectory
import base64

import fitz
import jieba
import matplotlib.pyplot as plt
import pandas as pd
import plotly.express as px
import pytesseract
import requests
import streamlit as st
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from matplotlib import font_manager
from pdf2image import convert_from_bytes
from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from pptx import Presentation
from wordcloud import WordCloud

from tools.file_router import route_file
from tools.prompt_templates import (
    prompt_evidence_contract,
    prompt_language_contract,
    prompt_material_label,
    summary_output_spec,
    top_venue_abstract_contract,
)


# 代码导航：
# 1. 配置与存储：Tesseract、SQLite、用户和历史记录。
# 2. 文本质量：乱码修复、OCR 清洗、可读性检查。
# 3. 文件解析：PDF、Word、PPT、图片、表格和文本。
# 4. 摘要生成：本地兜底、大模型提示词、PDCA、顶刊/顶会最终审稿。
# 5. 页面展示：分类、摘要双语分框、词频图、下载导出。
# 详细说明见 docs/代码结构与维护说明.md。

# 项目根目录和运行数据目录统一在这里配置，便于后续迁移、部署和排查问题。
BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
UPLOAD_DIR = DATA_DIR / "uploads"
DB_PATH = DATA_DIR / "app.db"
ACADEMIC_EXAMPLES_PATH = BASE_DIR / "academic_abstract_examples.jsonl"
REPORT_TEMPLATE_PATH = BASE_DIR / "templates" / "report_template.tex"

# 标准分类标签集中放在这里，侧边栏手动分类和历史筛选都复用同一套文档用途大类。
STANDARD_CATEGORY_RULES = {
    "学术论文/研究报告": ["摘要", "关键词", "引言", "方法", "实验", "结果", "结论", "参考文献", "研究", "模型", "数据集", "abstract", "method", "results", "references", "journal", "conference"],
    "商业分析/咨询报告": ["商业分析", "市场分析", "行业分析", "竞争", "战略", "客户", "用户", "增长", "收入", "成本", "利润", "商业模式", "咨询", "market", "strategy", "business"],
    "数据报表/统计分析": ["报表", "统计", "指标", "同比", "环比", "趋势", "看板", "销售额", "转化率", "均值", "方差", "样本", "dashboard", "report", "kpi"],
    "财务/审计/票据": ["财务", "审计", "会计", "发票", "金额", "税率", "付款", "收款", "银行", "费用", "报销", "凭证", "现金流", "资产负债率"],
    "年报/公告/披露文本": ["年度报告", "公告", "披露", "董事会报告", "管理层讨论", "重大事项", "经营情况", "公司治理", "募集说明书", "评级报告"],
    "法律合同/协议": ["合同", "协议", "甲方", "乙方", "违约责任", "条款", "履行", "交付", "保密", "争议解决", "法律", "诉讼"],
    "教学课件/课程材料": ["课程", "课件", "教学", "课堂", "学生", "教师", "学习目标", "练习", "作业", "考试", "知识点", "ppt", "lecture"],
    "政策通知/公文": ["通知", "公告", "会议", "决定", "要求", "单位", "日期", "请于", "安排", "报送", "政策", "文件精神"],
    "简历/人才材料": ["简历", "教育经历", "工作经历", "项目经验", "技能", "求职", "实习", "毕业院校", "候选人", "岗位"],
    "医学健康材料": ["医学", "临床", "疾病", "患者", "诊断", "治疗", "药物", "医院", "健康", "病例", "medical", "clinical"],
    "工程技术文档": ["工程", "技术方案", "系统设计", "架构", "算法", "接口", "部署", "测试", "性能", "代码", "计算机", "人工智能"],
    "图片/扫描资料": ["图片", "扫描", "截图", "页面", "图示", "海报", "表单", "ocr", "image", "scan"],
    "综合/通用材料": [],
}
STANDARD_CATEGORIES = list(STANDARD_CATEGORY_RULES.keys())


def configure_tesseract() -> None:
    """按优先级查找 Tesseract，并把可用路径写入 pytesseract 配置。"""
    candidates = [
        os.getenv("TESSERACT_CMD"),
        BASE_DIR / "tools" / "Tesseract-OCR" / "tesseract.exe",
        BASE_DIR / "tesseract" / "tesseract.exe",
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
        shutil.which("tesseract"),
    ]

    for candidate in candidates:
        if not candidate:
            continue
        command = Path(candidate)
        if command.exists():
            pytesseract.pytesseract.tesseract_cmd = str(command)
            tessdata_dir = command.parent / "tessdata"
            if tessdata_dir.exists() and not os.getenv("TESSDATA_PREFIX"):
                os.environ["TESSDATA_PREFIX"] = str(tessdata_dir)
            return


configure_tesseract()


def windows_ocr_engine_available() -> bool:
    """检查当前 Windows 环境是否可以调用系统内置 OCR 作为兜底。"""
    if os.name != "nt":
        return False
    script = BASE_DIR / "tools" / "windows_ocr.ps1"
    if not script.exists() or not shutil.which("powershell"):
        return False
    command = (
        "try {"
        "[Windows.Media.Ocr.OcrEngine, Windows.Foundation, ContentType=WindowsRuntime] | Out-Null;"
        "$engine=[Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages();"
        "if ($null -eq $engine) { exit 1 } else { exit 0 }"
        "} catch { exit 1 }"
    )
    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", command],
            capture_output=True,
            timeout=5,
        )
        return result.returncode == 0
    except Exception:
        return False


def get_tesseract_status() -> tuple[bool, str]:
    """返回 OCR 是否可用，以及展示给用户的 OCR 状态说明。"""
    command = getattr(pytesseract.pytesseract, "tesseract_cmd", "tesseract")
    command_path = Path(command)
    if command_path.exists() or shutil.which(command):
        return True, str(command)
    if windows_ocr_engine_available():
        return True, "未检测到 Tesseract，已启用 Windows 内置 OCR 兜底。图片 OCR 可用；扫描版 PDF 将使用 PyMuPDF 渲染后 OCR。"
    return False, (
        "未检测到 Tesseract OCR，且 Windows 内置 OCR 不可用。普通 PDF/Word 可正常解析；"
        "图片和扫描版 PDF 需要在本机安装 Tesseract，或将便携版放到 tools/Tesseract-OCR。"
        "云端部署时 packages.txt 会自动安装。"
    )

STOPWORDS = {
    "的", "了", "和", "是", "在", "有", "与", "及", "或", "为", "对", "等", "中",
    "根据", "以下", "论文", "内容", "生成", "中文", "学术", "摘要", "必要", "保留",
    "英文", "术语", "标题", "语言", "核心",
    "the", "and", "for", "with", "that", "this", "from", "are", "was", "were",
    "of", "in", "to", "as", "is", "on", "by", "an", "be", "we", "our", "it"
}


# ---------- 存储与用户认证 ----------

def init_storage() -> None:
    """初始化 SQLite 数据库、上传目录和必要的数据表。"""
    DATA_DIR.mkdir(exist_ok=True)
    UPLOAD_DIR.mkdir(exist_ok=True)
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT NOT NULL UNIQUE,
                password_hash TEXT NOT NULL,
                salt TEXT NOT NULL,
                created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER NOT NULL,
                filename TEXT NOT NULL,
                file_path TEXT NOT NULL,
                category TEXT NOT NULL,
                summary TEXT NOT NULL,
                pdca_report TEXT NOT NULL DEFAULT '',
                extracted_text TEXT NOT NULL,
                word_count INTEGER NOT NULL,
                processing_seconds REAL NOT NULL DEFAULT 0,
                created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY(user_id) REFERENCES users(id)
            )
            """
        )
        columns = {row[1] for row in conn.execute("PRAGMA table_info(documents)").fetchall()}
        if "pdca_report" not in columns:
            conn.execute("ALTER TABLE documents ADD COLUMN pdca_report TEXT NOT NULL DEFAULT ''")
        if "processing_seconds" not in columns:
            conn.execute("ALTER TABLE documents ADD COLUMN processing_seconds REAL NOT NULL DEFAULT 0")
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS parse_cache (
                cache_key TEXT PRIMARY KEY,
                file_hash TEXT NOT NULL,
                filename TEXT NOT NULL,
                parser_options TEXT NOT NULL,
                extracted_text TEXT NOT NULL,
                created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
            )
            """
        )


def hash_password(password: str, salt: str | None = None) -> tuple[str, str]:
    """使用 PBKDF2 对密码加盐哈希，避免明文保存用户密码。"""
    salt = salt or secrets.token_hex(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt.encode("utf-8"), 120_000)
    return digest.hex(), salt


def create_user(username: str, password: str) -> tuple[bool, str]:
    """创建新用户，并返回是否成功及页面提示信息。"""
    username = username.strip()
    if not re.fullmatch(r"[A-Za-z0-9_\u4e00-\u9fff]{2,20}", username):
        return False, "用户名需为 2-20 位中文、字母、数字或下划线。"
    if len(password) < 6:
        return False, "密码至少需要 6 位。"

    password_hash, salt = hash_password(password)
    try:
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute(
                "INSERT INTO users (username, password_hash, salt) VALUES (?, ?, ?)",
                (username, password_hash, salt),
            )
        return True, "注册成功，请登录。"
    except sqlite3.IntegrityError:
        return False, "用户名已存在。"


def authenticate(username: str, password: str) -> dict | None:
    """校验用户名和密码，成功时返回当前用户的基础信息。"""
    with sqlite3.connect(DB_PATH) as conn:
        conn.row_factory = sqlite3.Row
        user = conn.execute("SELECT * FROM users WHERE username = ?", (username.strip(),)).fetchone()
    if user is None:
        return None
    password_hash, _ = hash_password(password, user["salt"])
    if not secrets.compare_digest(password_hash, user["password_hash"]):
        return None
    return {"id": user["id"], "username": user["username"]}


def get_or_create_demo_user() -> dict:
    """为演示场景创建或读取一个免注册试用账号。"""
    username = "demo_user"
    password = "demo_password_123"
    user = authenticate(username, password)
    if user is not None:
        return user
    ok, _ = create_user(username, password)
    if ok:
        user = authenticate(username, password)
        if user is not None:
            return user
    with sqlite3.connect(DB_PATH) as conn:
        conn.row_factory = sqlite3.Row
        row = conn.execute("SELECT id, username FROM users WHERE username = ?", (username,)).fetchone()
    if row is None:
        raise RuntimeError("免注册试用账号创建失败。")
    return {"id": row["id"], "username": row["username"]}


def is_demo_user(user: dict | None) -> bool:
    """判断当前账号是否为免注册试用账号。"""
    return bool(user and user.get("username") == "demo_user")


def delete_demo_documents() -> None:
    """清空免注册试用账号产生的历史记录和上传文件。"""
    with sqlite3.connect(DB_PATH) as conn:
        conn.row_factory = sqlite3.Row
        demo = conn.execute("SELECT id FROM users WHERE username = ?", ("demo_user",)).fetchone()
        if demo is None:
            return
        rows = conn.execute("SELECT file_path FROM documents WHERE user_id = ?", (demo["id"],)).fetchall()
        conn.execute("DELETE FROM documents WHERE user_id = ?", (demo["id"],))
    for row in rows:
        path = Path(row["file_path"])
        try:
            if path.exists() and UPLOAD_DIR in path.resolve().parents:
                path.unlink()
        except Exception:
            pass
    demo_dir = UPLOAD_DIR / str(demo["id"])
    try:
        if demo_dir.exists() and not any(demo_dir.iterdir()):
            demo_dir.rmdir()
    except Exception:
        pass


def save_document(user_id: int, filename: str, file_bytes: bytes, category: str, summary: str, pdca_report: str, extracted_text: str, word_count: int, processing_seconds: float = 0.0) -> None:
    """保存用户上传的原始文件和对应的分析结果。"""
    user_dir = UPLOAD_DIR / str(user_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^A-Za-z0-9_.\-\u4e00-\u9fff]", "_", filename)
    stored_name = f"{int(time.time())}_{uuid.uuid4().hex}_{safe_name}"
    file_path = user_dir / stored_name
    file_path.write_bytes(file_bytes)

    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            INSERT INTO documents
            (user_id, filename, file_path, category, summary, pdca_report, extracted_text, word_count, processing_seconds)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (user_id, filename, str(file_path), category, summary, pdca_report, extracted_text, word_count, float(processing_seconds or 0.0)),
        )


def list_user_documents(user_id: int) -> list[sqlite3.Row]:
    """读取当前用户自己的文档历史记录。"""
    with sqlite3.connect(DB_PATH) as conn:
        conn.row_factory = sqlite3.Row
        return conn.execute(
            """
            SELECT id, filename, category, summary, word_count, processing_seconds, created_at
            FROM documents
            WHERE user_id = ?
            ORDER BY created_at DESC, id DESC
            """,
            (user_id,),
        ).fetchall()


def get_user_document(user_id: int, document_id: int) -> sqlite3.Row | None:
    """按用户和文档编号读取单条历史记录，防止跨用户访问。"""
    with sqlite3.connect(DB_PATH) as conn:
        conn.row_factory = sqlite3.Row
        return conn.execute(
            """
            SELECT * FROM documents
            WHERE id = ? AND user_id = ?
            """,
            (document_id, user_id),
        ).fetchone()


def build_parse_cache_key(file_hash: str, filename: str, ocr_mode: str, pdf_extract_mode: str) -> str:
    """按文件内容和解析参数生成缓存键。"""
    suffix = Path(filename).suffix.lower()
    raw = f"{file_hash}:{suffix}:{ocr_mode}:{pdf_extract_mode}:v2"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def get_parse_cache(cache_key: str) -> str | None:
    """读取解析缓存，避免同一文件重复 OCR 或重复调用多模态接口。"""
    with sqlite3.connect(DB_PATH) as conn:
        row = conn.execute("SELECT extracted_text FROM parse_cache WHERE cache_key = ?", (cache_key,)).fetchone()
    return row[0] if row else None


def save_parse_cache(cache_key: str, file_hash: str, filename: str, parser_options: str, extracted_text: str) -> None:
    """保存解析缓存。"""
    if not clean_text(extracted_text):
        return
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """
            INSERT OR REPLACE INTO parse_cache
            (cache_key, file_hash, filename, parser_options, extracted_text)
            VALUES (?, ?, ?, ?, ?)
            """,
            (cache_key, file_hash, filename, parser_options, extracted_text),
        )


def detect_file_kind(file_bytes: bytes, filename: str) -> str:
    """基于文件头识别真实格式，降低扩展名错误导致的解析失败。"""
    suffix = Path(filename).suffix.lower().lstrip(".")
    head = file_bytes[:16]
    if head.startswith(b"%PDF"):
        return "pdf"
    if head.startswith(b"PK\x03\x04"):
        if suffix in {"docx", "pptx", "xlsx"}:
            return suffix
        return "zip-office"
    if head.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        if suffix in {"doc", "ppt", "xls"}:
            return suffix
        return "ole-office"
    if head.startswith(b"\x89PNG"):
        return "png"
    if head.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if head.startswith((b"II*\x00", b"MM\x00*")):
        return "tiff"
    if head.startswith(b"BM"):
        return "bmp"
    return suffix or "unknown"


def validate_file_kind(file_bytes: bytes, filename: str) -> None:
    """检查扩展名和真实文件头是否明显不一致。"""
    suffix = Path(filename).suffix.lower().lstrip(".")
    kind = detect_file_kind(file_bytes, filename)
    if suffix in {"txt", "md", "json", "jsonl", "csv"}:
        return
    compatible = {
        "jpg": {"jpg", "jpeg"},
        "jpeg": {"jpg", "jpeg"},
        "tif": {"tif", "tiff"},
        "tiff": {"tif", "tiff"},
        "zip-office": {"docx", "pptx", "xlsx"},
        "ole-office": {"doc", "ppt", "xls"},
    }
    allowed = compatible.get(kind, {kind})
    if suffix and suffix not in allowed:
        raise ValueError(f"文件扩展名 .{suffix} 与真实格式 {kind} 不一致，请另存为正确格式后重新上传。")


# ---------- 文档解析与 OCR ----------

def mojibake_marker_count(text: str) -> int:
    """统计常见中文编码乱码标记数量，用于判断是否需要自动修复。"""
    marker_codes = (
        0x9418, 0x9354, 0x7035, 0x951B, 0x9286, 0x7ECB, 0x6D93, 0x93C2,
        0x95C8, 0x93C8, 0x7455, 0x71AC, 0xFFFD, 0x9428, 0x93C2, 0x6D60,
        0x9366, 0x5997, 0x704F, 0x20AC,
    )
    markers = [chr(code) for code in marker_codes]
    return sum(text.count(marker) for marker in markers)


def text_readability_score(text: str) -> float:
    """给文本可读性打分，分数越高表示越像正常中英文文本。"""
    if not text:
        return 0.0
    compact = re.sub(r"\s+", "", text)
    if not compact:
        return 0.0
    readable = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", compact))
    noise = mojibake_marker_count(compact) * 8 + compact.count("?") + compact.count(chr(0xFFFD)) * 10
    return readable / max(len(compact), 1) - noise / max(len(compact), 1)


def repair_mojibake_text(text: str) -> str:
    """尝试修复 UTF-8 被误按 GBK/ANSI 读取造成的中文乱码。"""
    if not text:
        return ""
    original = text
    candidates = [original]
    for source_encoding in ("gbk", "gb18030", "latin1"):
        try:
            candidate = original.encode(source_encoding, errors="strict").decode("utf-8", errors="strict")
        except (UnicodeEncodeError, UnicodeDecodeError):
            candidate = ""
        if candidate:
            candidates.append(candidate)
        try:
            candidate = original.encode(source_encoding, errors="ignore").decode("utf-8", errors="ignore")
        except UnicodeDecodeError:
            candidate = ""
        if candidate:
            candidates.append(candidate)

    best = max(candidates, key=text_readability_score)
    original_score = text_readability_score(original)
    best_score = text_readability_score(best)
    if best != original and best_score > original_score + 0.08:
        return best
    return original


def clean_text(text: str) -> str:
    """清理通用文本中的空字符、多余空格和过多换行。"""
    text = repair_mojibake_text(text)
    text = unicodedata.normalize("NFKC", text)
    formula_replacements = {
        "≤": "<=",
        "≥": ">=",
        "≠": "!=",
        "×": r"\times ",
        "÷": r"\div ",
        "∑": r"\sum ",
        "∏": r"\prod ",
        "√": r"\sqrt ",
        "∞": r"\infty ",
        "≈": r"\approx ",
        "→": r"\to ",
        "←": r"\leftarrow ",
        "∈": r"\in ",
        "∂": r"\partial ",
        "α": r"\alpha ",
        "β": r"\beta ",
        "γ": r"\gamma ",
        "δ": r"\delta ",
        "λ": r"\lambda ",
        "μ": r"\mu ",
        "σ": r"\sigma ",
        "θ": r"\theta ",
        "ρ": r"\rho ",
    }
    for source, target in formula_replacements.items():
        text = text.replace(source, target)
    text = text.replace("\x00", " ")
    text = re.sub(r"[\u200b\u200c\u200d\ufeff]", "", text)
    text = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f]", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def looks_like_formula_line(line: str) -> bool:
    """识别公式或数学表达式行，避免清洗和断行合并时破坏结构。"""
    stripped = line.strip()
    if len(stripped) < 3:
        return False
    formula_marks = r"=+\-*/^_{}\\∑∏√∞≈≤≥≠<>∂∫αβγδλμσθρ"
    mark_count = len(re.findall(f"[{re.escape(formula_marks)}]", stripped))
    latin_math = len(re.findall(r"\b[A-Za-z]\w*\b", stripped))
    if "$" in stripped or r"\frac" in stripped or r"\sum" in stripped or r"\int" in stripped:
        return True
    if "=" in stripped and mark_count >= 2:
        return True
    return mark_count >= 4 and latin_math >= 1


def line_signature(line: str) -> str:
    """生成行文本指纹，用于去掉重复页眉、页脚或连续重复行。"""
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "", line).lower()


def should_preserve_linebreak(previous: str, current: str, source_type: str) -> bool:
    """判断两行之间是否应保留换行，避免把标题、列表、表格和 PPT 页码合并进正文。"""
    if not previous or not current:
        return True
    if looks_like_formula_line(previous) or looks_like_formula_line(current):
        return True
    if source_type in {"ppt", "table"}:
        return True
    if "|" in previous or "|" in current:
        return True
    if re.match(r"^(第\d+页|Page\s+\d+|Slide\s+\d+|备注[:：])", previous, re.IGNORECASE):
        return True
    if re.match(r"^(第\d+页|Page\s+\d+|Slide\s+\d+|备注[:：])", current, re.IGNORECASE):
        return True
    if re.match(r"^([一二三四五六七八九十]+[、.]|\d+[、.]|\(\d+\)|[A-Za-z]\.)", current):
        return True
    if previous.endswith(("。", "！", "？", ".", "!", "?", "；", ";", "：", ":")):
        return True
    if len(previous) <= 18 and not previous.endswith(("，", ",")):
        return True
    return False


def join_extracted_lines(lines: list[str], source_type: str) -> str:
    """合并 PDF/Word 提取中被硬换行切断的正文行，同时尽量保留结构。"""
    merged: list[str] = []
    for line in lines:
        line = line.strip()
        if not line:
            if merged and merged[-1] != "":
                merged.append("")
            continue
        if not merged or merged[-1] == "" or should_preserve_linebreak(merged[-1], line, source_type):
            merged.append(line)
            continue
        separator = "" if re.search(r"[\u4e00-\u9fff]$", merged[-1]) and re.match(r"^[\u4e00-\u9fff]", line) else " "
        merged[-1] = f"{merged[-1].rstrip('- ')}{separator}{line}"
    return "\n".join(merged)


def normalize_table_cell(value) -> str:
    """规范表格单元格文本，保留空单元格占位，避免列错位。"""
    text = clean_text("" if value is None else str(value))
    return text.replace("|", "\\|").replace("\n", " ")


def markdown_table_from_rows(rows: list[list], max_rows: int | None = None) -> str:
    """把二维数组转成 Markdown 表格，便于表格在解析、摘要和导出中保持结构。"""
    if max_rows is not None:
        rows = rows[:max_rows]
    normalized = [[normalize_table_cell(cell) for cell in row] for row in rows]
    normalized = [row for row in normalized if any(cell.strip() for cell in row)]
    if not normalized:
        return ""
    width = max(len(row) for row in normalized)
    normalized = [row + [""] * (width - len(row)) for row in normalized]
    header = normalized[0]
    body = normalized[1:] or [[""] * width]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def post_process_extracted_text(text: str, source_type: str = "document") -> str:
    """对全文提取结果做结构化清洗：去重复、修断行、保留段落。"""
    text = clean_text(text)
    if not text:
        return ""

    raw_lines = [line.strip() for line in text.splitlines()]
    signature_counts = Counter(
        line_signature(line)
        for line in raw_lines
        if 2 <= len(line_signature(line)) <= 40
    )
    cleaned_lines = []
    previous_signature = ""
    for line in raw_lines:
        signature = line_signature(line)
        if not line:
            cleaned_lines.append("")
            previous_signature = ""
            continue
        if signature and signature == previous_signature:
            continue
        if source_type == "pdf" and signature_counts.get(signature, 0) >= 3 and len(line) <= 45:
            continue
        cleaned_lines.append(line)
        previous_signature = signature

    text = join_extracted_lines(cleaned_lines, source_type)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def clean_ocr_text(text: str) -> str:
    """清理 OCR 结果中的常见识别错误、异常空格和标点空白。"""
    text = clean_text(text)
    if not text:
        return ""

    replacements = {
        "竟争": "竞争",
        "哽快": "更快",
        "哽深": "更深",
        "昊": "吴",
        "賣化": "变化",
        "慌秀": "演示",
        "趼": "研",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)

    # Windows 内置 OCR 常把中文字符之间插入空格，这里合并回正常中文句子。
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+([，。！？；：、])", r"\1", text)
    text = re.sub(r"([，。！？；：、])\s+(?=[\u4e00-\u9fff])", r"\1", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()


def ocr_text_score(text: str) -> float:
    """给 OCR 候选结果打分，用于从多种图像预处理结果中选最可信文本。"""
    if not text:
        return 0.0
    compact = re.sub(r"\s+", "", text)
    cjk = len(re.findall(r"[\u4e00-\u9fff]", compact))
    latin_digits = len(re.findall(r"[A-Za-z0-9]", compact))
    replacement_char = chr(0xFFFD)
    punctuation_noise = len(re.findall(f"[□■●◆{replacement_char}#{{}}|~]", compact))
    spaced_cjk = len(re.findall(r"[\u4e00-\u9fff]\s+[\u4e00-\u9fff]", text))
    return cjk * 2.0 + latin_digits * 0.8 - punctuation_noise * 3.0 - spaced_cjk * 0.8


def ocr_image_variants(image: Image.Image, mode: str = "standard") -> list[Image.Image]:
    """按 OCR 模式生成图像预处理版本，平衡识别速度和识别质量。"""
    base = image.convert("RGB")
    variants = [base]
    if mode == "fast":
        return variants

    width, height = base.size
    scale = 2 if max(width, height) < 1800 else 1
    enlarged = base.resize((width * scale, height * scale), Image.Resampling.LANCZOS) if scale > 1 else base
    gray = ImageOps.grayscale(enlarged)
    contrast = ImageEnhance.Contrast(gray).enhance(1.8)
    sharp = contrast.filter(ImageFilter.SHARPEN)
    threshold = sharp.point(lambda pixel: 255 if pixel > 175 else 0)
    variants.extend([sharp.convert("RGB")])
    if mode == "accurate":
        variants.extend([enlarged, threshold.convert("RGB")])
    return variants


def is_unreadable_text(text: str) -> bool:
    """判断提取文本是否疑似乱码或不可读，避免把坏文本继续送去摘要。"""
    compact = re.sub(r"\s+", "", text)
    if len(compact) < 30:
        return False
    readable_chars = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", compact))
    readable_ratio = readable_chars / len(compact)
    question_ratio = compact.count("?") / len(compact)
    replacement_char = chr(0xFFFD)
    replacement_ratio = compact.count(replacement_char) / len(compact)
    mojibake_chars = {replacement_char, "□", "■", "●", "◆", "◇"}
    mojibake_ratio = (sum(1 for char in compact if char in mojibake_chars) + mojibake_marker_count(compact)) / len(compact)
    if replacement_ratio > 0.03 or mojibake_ratio > 0.08 or looks_mojibake(compact):
        return True
    return readable_ratio < 0.15 and question_ratio > 0.85


def split_text_chunks(text: str, chunk_size: int = 5000, overlap: int = 300) -> list[str]:
    """把长文本切成带重叠的片段，供大模型分块摘要使用。"""
    text = text.strip()
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks



def load_academic_examples(limit: int = 200) -> list[dict[str, str]]:
    """加载本地学术摘要范例库，用于提示词中的风格参考。"""
    if not ACADEMIC_EXAMPLES_PATH.exists():
        return []
    examples = []
    for line in ACADEMIC_EXAMPLES_PATH.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            item = json.loads(line)
        except json.JSONDecodeError:
            continue
        abstract = str(item.get("abstract", "")).strip()
        source = str(item.get("source", "")).strip()
        title = str(item.get("title", "")).strip()
        if abstract and not looks_mojibake(abstract):
            examples.append({"title": title, "source": source, "abstract": abstract})
        if len(examples) >= limit:
            break
    return examples


def looks_mojibake(text: str) -> bool:
    """识别常见中文编码乱码，避免坏样例或坏输出污染摘要。"""
    if not text:
        return False
    marker_count = mojibake_marker_count(text)
    return marker_count >= 3 and marker_count / max(len(text), 1) > 0.01


def lexical_terms(text: str) -> set[str]:
    """提取中英文关键词集合，用于匹配相近的学术摘要范例。"""
    latin = re.findall(r"[a-zA-Z][a-zA-Z0-9_-]{2,}", text.lower())
    chinese = re.findall(r"[\u4e00-\u9fff]{2,}", text)
    return set(latin + chinese)


def academic_example_domain_score(query_text: str, example: dict[str, str]) -> int:
    """根据文档主题给摘要范例额外加权，让风格样例更贴近当前材料。"""
    haystack = f"{example.get('title', '')} {example.get('abstract', '')}"
    domain_rules = [
        (["财务困境", "风险预警", "违约", "破产", "st", "catboost", "排序"], ["财务困境", "风险排序", "预警"]),
        (["公司治理", "董事会", "股东", "中小股东", "代理成本", "治理机制"], ["公司治理", "董事会", "股东"]),
        (["streamlit", "ocr", "pdf", "word", "系统", "部署", "文档处理"], ["系统", "文档智能", "ocr"]),
        (["附录", "表a", "表b", "supplement", "appendix", "稳健性"], ["补充材料", "附加表格", "稳健性"]),
        (["literature", "review", "综述", "文献", "研究现状"], ["综述", "文献", "知识结构"]),
        (["transformer", "machine learning", "模型", "算法", "dataset", "benchmark"], ["机器学习", "模型", "实验"]),
    ]
    score = 0
    query_lower = query_text.lower()
    haystack_lower = haystack.lower()
    for query_keywords, example_keywords in domain_rules:
        query_hits = sum(1 for keyword in query_keywords if keyword.lower() in query_lower)
        example_hits = sum(1 for keyword in example_keywords if keyword.lower() in haystack_lower)
        if query_hits and example_hits:
            score += query_hits * example_hits * 8
    return score


def select_academic_examples(text: str, max_examples: int = 3) -> list[dict[str, str]]:
    """根据当前文档词汇重叠度选择最相关的摘要范例。"""
    examples = load_academic_examples()
    if not examples:
        return []
    query_terms = lexical_terms(text)
    scored = []
    for item in examples:
        haystack = " ".join([item.get("title", ""), item.get("source", ""), item.get("abstract", "")])
        overlap = len(query_terms & lexical_terms(haystack))
        domain_score = academic_example_domain_score(text, item)
        scored.append((overlap + domain_score, item))
    scored.sort(key=lambda pair: pair[0], reverse=True)
    return [item for score, item in scored[:max_examples] if score > 0] or [item for _, item in scored[:max_examples]]


def format_academic_examples(text: str) -> str:
    """把选中的摘要范例格式化为可直接放入大模型提示词的文本。"""
    examples = select_academic_examples(text)
    if not examples:
        return "No reference abstracts are available."
    parts = []
    for index, item in enumerate(examples, start=1):
        abstract = item["abstract"][:900]
        title = item.get("title") or f"Example {index}"
        parts.append(f"范例{index}标题：{title}\n范例{index}摘要风格参考：\n{abstract}")
    return "\n\n".join(parts)


def run_ocr(image: Image.Image, mode: str = "standard") -> str:
    """优先使用 Tesseract OCR，并从多种预处理和版面参数中选择最佳结果。"""
    try:
        best_text = ""
        best_score = 0.0
        configs = ["--oem 1 --psm 6 -c preserve_interword_spaces=1"]
        if mode != "fast":
            configs.append("--oem 1 --psm 3 -c preserve_interword_spaces=1")
            configs.append("--oem 1 --psm 11 -c preserve_interword_spaces=1")
        if mode == "accurate":
            configs.append("--oem 1 --psm 12 -c preserve_interword_spaces=1")
        for variant in ocr_image_variants(image, mode=mode):
            for config in configs:
                text = clean_ocr_text(pytesseract.image_to_string(variant, lang="chi_sim+eng", config=config))
                score = ocr_text_score(text)
                if score > best_score:
                    best_text = text
                    best_score = score
        return best_text
    except pytesseract.TesseractNotFoundError as exc:
        fallback_text = run_windows_ocr(image, mode=mode)
        if fallback_text:
            return fallback_text
        raise RuntimeError(
            "未找到 Tesseract OCR，且 Windows 内置 OCR 也不可用。本地测试图片/扫描版 PDF 时，请运行 "
            "install_ocr_windows.ps1，或把便携版 Tesseract 放到 tools/Tesseract-OCR/tesseract.exe，"
            "并包含 tessdata/chi_sim.traineddata。普通文字 PDF 和 Word 不需要 OCR；"
            "部署到 Streamlit Cloud 时会根据 packages.txt 自动安装 OCR。"
        ) from exc
    except pytesseract.TesseractError as exc:
        message = str(exc)
        if "chi_sim" in message or "Failed loading language" in message:
            fallback_text = run_windows_ocr(image, mode=mode)
            if fallback_text:
                return fallback_text
            raise RuntimeError(
                "Tesseract 已找到，但缺少中文语言包 chi_sim。请确保 tools/Tesseract-OCR/tessdata/chi_sim.traineddata 存在，或安装中文简体语言数据。"
            ) from exc
        raise RuntimeError(f"OCR 识别失败：{message}") from exc


def run_windows_ocr(image: Image.Image, mode: str = "standard") -> str:
    """调用 Windows OCR 脚本，并按模式从预处理图片中选择最佳识别结果。"""
    if os.name != "nt":
        return ""
    script = BASE_DIR / "tools" / "windows_ocr.ps1"
    if not script.exists() or not shutil.which("powershell"):
        return ""
    best_text = ""
    best_score = 0.0
    temp_paths = []
    try:
        for variant in ocr_image_variants(image, mode=mode):
            with NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                temp_path = temp_file.name
                temp_paths.append(temp_path)
            variant.convert("RGB").save(temp_path, format="PNG")
            result = subprocess.run(
                [
                    "powershell",
                    "-NoProfile",
                    "-ExecutionPolicy",
                    "Bypass",
                    "-File",
                    str(script),
                    "-ImagePath",
                    temp_path,
                ],
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=30,
            )
            if result.returncode != 0:
                continue
            text = clean_ocr_text(result.stdout)
            score = ocr_text_score(text)
            if score > best_score:
                best_text = text
                best_score = score
        return best_text
    except Exception:
        return ""
    finally:
        for temp_path in temp_paths:
            try:
                Path(temp_path).unlink(missing_ok=True)
            except Exception:
                pass


def render_pdf_pages_with_fitz(file_bytes: bytes, mode: str = "standard", max_pages: int | None = None) -> list[Image.Image]:
    """用 PyMuPDF 将 PDF 页面渲染成图片，作为 Poppler 失败时的兜底方案。"""
    images = []
    zoom = {"fast": 1.4, "standard": 1.7, "accurate": 2.2}.get(mode, 1.7)
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page_count = len(doc) if max_pages is None else min(len(doc), max_pages)
    for page_index in range(page_count):
        page = doc[page_index]
        pixmap = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        images.append(Image.open(io.BytesIO(pixmap.tobytes("png"))))
    return images


def render_pdf_page_with_fitz(doc: fitz.Document, page_index: int, mode: str = "standard") -> Image.Image:
    """把指定 PDF 页渲染为图片，用于坏页多模态补救。"""
    zoom = {"fast": 1.4, "standard": 1.7, "accurate": 2.2}.get(mode, 1.7)
    page = doc[page_index]
    pixmap = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    return Image.open(io.BytesIO(pixmap.tobytes("png")))


def image_to_data_url(image: Image.Image, max_width: int = 1400) -> str:
    """把页面截图压缩为 data URL，供 OpenAI 兼容多模态接口使用。"""
    image = image.convert("RGB")
    if image.width > max_width:
        ratio = max_width / image.width
        image = image.resize((max_width, max(1, int(image.height * ratio))), Image.Resampling.LANCZOS)
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", quality=82, optimize=True)
    payload = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/jpeg;base64,{payload}"


def extract_pdf_page_text(page: fitz.Page) -> str:
    """同时尝试普通文本流和坐标块文本，选择更可读的一版。"""
    plain_text = clean_text(page.get_text("text"))
    table_text = ""
    try:
        table_finder = page.find_tables()
        tables = []
        for table_index, table in enumerate(getattr(table_finder, "tables", []), start=1):
            rows = table.extract()
            table_markdown = markdown_table_from_rows(rows)
            if table_markdown:
                tables.append(f"PDF表格{table_index}：\n{table_markdown}")
        table_text = clean_text("\n\n".join(tables))
    except Exception:
        table_text = ""
    try:
        blocks = page.get_text("blocks")
        text_blocks = []
        for block in sorted(blocks, key=lambda item: (round(item[1] / 12), item[0])):
            block_text = clean_text(str(block[4])) if len(block) > 4 else ""
            if block_text:
                text_blocks.append(block_text)
        block_text = clean_text("\n".join(text_blocks))
    except Exception:
        block_text = ""

    candidates = [item for item in [plain_text, block_text] if item]
    if not candidates:
        return table_text

    def candidate_score(text: str) -> float:
        compact = re.sub(r"\s+", "", text)
        structure_bonus = min(text.count("\n"), 30) * 0.4
        length_bonus = min(len(compact), 2000) / 200
        return text_readability_score(text) * 100 + structure_bonus + length_bonus

    best_text = max(candidates, key=candidate_score)
    if table_text and table_text not in best_text:
        return clean_text(f"{best_text}\n\n{table_text}")
    return best_text


def pdf_text_needs_vision(text: str) -> bool:
    """判断 PDF 本地文本是否明显不足，需要尝试多模态页面读取。"""
    cleaned = clean_text(text)
    compact = re.sub(r"\s+", "", cleaned)
    if len(compact) < 220:
        return True
    if is_unreadable_text(cleaned):
        return True
    formula_lines = [line for line in cleaned.splitlines() if looks_like_formula_line(line)]
    if formula_lines and len(formula_lines) / max(len([line for line in cleaned.splitlines() if line.strip()]), 1) > 0.18:
        return True
    short_lines = [line for line in cleaned.splitlines() if 1 <= len(line.strip()) <= 8]
    lines = [line for line in cleaned.splitlines() if line.strip()]
    return bool(lines) and len(short_lines) / max(len(lines), 1) > 0.55


def call_vision_pdf_extraction(images: list[Image.Image], filename: str = "", page_numbers: list[int] | None = None) -> str:
    """调用多模态模型从 PDF 页面截图中提取正文和表格要点。"""
    api_key, api_base, model = get_vision_llm_config()
    if not api_key:
        raise RuntimeError("未配置视觉模型 API。")
    content = [
        {
            "type": "text",
            "text": (
                "请从这些PDF页面截图中提取可读正文。要求：\n"
                "1. 按页面顺序输出，保留标题、段落、列表、图片说明和表格关键信息。\n"
                "2. 不要总结，不要改写，不要补充原图没有的信息。\n"
                "3. 对数学公式、计量模型、约束条件和变量定义，尽量用 LaTeX 原样转写，例如 y_{i,t}=\\alpha+\\beta x_{i,t}+\\epsilon_{i,t}。\n"
                "4. 公式单独成行，变量下标、上标、求和、分式、希腊字母和比较符号尽量保留；不要把公式改写成文字解释。\n"
                "5. 表格必须转写为 Markdown 表格，保留列名、行名、空单元格和数值单位；不要只概括表格。\n"
                "6. 图片、图表或流程图请保留图题、坐标轴、图例、节点关系和图中文字；读不清的数值写[无法识别]。\n"
                "7. 识别不清的地方写[无法识别]，不要猜测。\n"
                f"文件名：{filename or '未命名PDF'}"
            ),
        }
    ]
    page_numbers = page_numbers or list(range(1, len(images) + 1))
    for index, image in enumerate(images, start=1):
        page_number = page_numbers[index - 1] if index - 1 < len(page_numbers) else index
        content.append({"type": "text", "text": f"第{page_number}页："})
        content.append({"type": "image_url", "image_url": {"url": image_to_data_url(image)}})
    response = requests.post(
        f"{api_base}/chat/completions",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json={
            "model": model,
            "messages": [
                {"role": "system", "content": "你是严谨的PDF页面文字提取助手，只做OCR、公式转写和结构化转写，不做摘要。"},
                {"role": "user", "content": content},
            ],
            "temperature": 0,
            "max_tokens": 2200,
            "stream": True,
        },
        stream=True,
        timeout=90,
    )
    response.raise_for_status()
    return clean_text(parse_llm_response(response))


def call_vision_image_extraction(image: Image.Image) -> str:
    """用多模态模型结构化识别单张图片、图表、流程图或截图。"""
    api_key, api_base, model = get_vision_llm_config()
    if not api_key:
        raise RuntimeError("未配置视觉模型 API。")
    response = requests.post(
        f"{api_base}/chat/completions",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json={
            "model": model,
            "messages": [
                {"role": "system", "content": "你是严谨的图片文字、图表和表格转写助手，只做结构化识别，不做摘要。"},
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "请结构化转写这张图片中的信息：\n"
                                "1. 提取所有可读文字，保持原顺序。\n"
                                "2. 若有表格，转为 Markdown 表格，保留列名、行名、空单元格、数值和单位。\n"
                                "3. 若有图表，写出图题、坐标轴、图例、系列名称、关键数值和趋势，不要凭空补数。\n"
                                "4. 若有流程图，按节点和箭头关系列出。\n"
                                "5. 若有公式，尽量用 LaTeX 原样转写。\n"
                                "6. 识别不清处写[无法识别]。"
                            ),
                        },
                        {"type": "image_url", "image_url": {"url": image_to_data_url(image)}},
                    ],
                },
            ],
            "temperature": 0,
            "max_tokens": 1800,
            "stream": True,
        },
        stream=True,
        timeout=90,
    )
    response.raise_for_status()
    return clean_text(parse_llm_response(response))


def split_extracted_pages(text: str) -> dict[int, str]:
    """把多模态返回的按页文本拆成 {页码: 正文}，便于替换坏页。"""
    text = clean_text(text)
    if not text:
        return {}
    matches = list(re.finditer(r"(?:^|\n)\s*第\s*(\d+)\s*页\s*[:：]?\s*", text))
    if not matches:
        return {1: text}
    pages: dict[int, str] = {}
    for index, match in enumerate(matches):
        page_number = int(match.group(1))
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        content = clean_text(text[start:end])
        if content:
            pages[page_number] = content
    return pages


def extract_pdf(file_bytes: bytes, ocr_mode: str = "standard", pdf_extract_mode: str = "auto", filename: str = "") -> str:
    """逐页提取 PDF：好页用文本，坏页才用多模态或 OCR 兜底。"""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page_entries: list[dict] = []
    weak_entries: list[dict] = []
    for index, page in enumerate(doc, start=1):
        page_text = extract_pdf_page_text(page)
        weak = pdf_extract_mode == "vision" or (pdf_extract_mode == "auto" and pdf_text_needs_vision(page_text))
        entry = {"page": index, "text": page_text, "weak": weak}
        page_entries.append(entry)
        if weak:
            weak_entries.append(entry)

    if pdf_extract_mode == "text":
        return clean_text("\n\n".join(f"第{entry['page']}页\n{entry['text']}" for entry in page_entries if entry["text"]))

    max_vision_pages = int(get_secret("VISION_PDF_MAX_PAGES") or 6)
    if weak_entries and pdf_extract_mode in {"auto", "vision"}:
        selected = weak_entries[:max_vision_pages]
        try:
            images = [render_pdf_page_with_fitz(doc, entry["page"] - 1, mode="standard") for entry in selected]
            page_numbers = [entry["page"] for entry in selected]
            vision_text = call_vision_pdf_extraction(images, filename=filename, page_numbers=page_numbers)
            if clean_text(vision_text):
                vision_by_page = split_extracted_pages(vision_text)
                for entry in selected:
                    replacement = clean_text(vision_by_page.get(entry["page"], ""))
                    if replacement and len(replacement) >= max(30, len(entry["text"]) * 0.5):
                        entry["text"] = replacement
        except Exception:
            if pdf_extract_mode == "vision":
                raise

    if any(entry["text"] and not pdf_text_needs_vision(entry["text"]) for entry in page_entries):
        return clean_text("\n\n".join(f"第{entry['page']}页\n{entry['text']}" for entry in page_entries if entry["text"]))

    try:
        dpi = {"fast": 130, "standard": 170, "accurate": 220}.get(ocr_mode, 170)
        images = convert_from_bytes(file_bytes, dpi=dpi)
    except Exception as exc:
        try:
            images = render_pdf_pages_with_fitz(file_bytes, mode=ocr_mode)
        except Exception as render_exc:
            raise RuntimeError(
                "扫描版 PDF 转图失败。"
                "云端请确认 packages.txt 包含 poppler-utils；"
                "Windows 本地可安装 Poppler，或使用 PyMuPDF 渲染兜底。"
            ) from render_exc

    ocr_text = []
    for index, image in enumerate(images, start=1):
        page_text = run_ocr(image, mode=ocr_mode)
        if page_text:
            ocr_text.append(f"第{index}页\n{page_text}")
    return clean_text("\n\n".join(ocr_text))


def extract_docx(file_obj) -> str:
    """提取 Word 文档中的段落和表格文本。"""
    document = Document(file_obj)
    parts = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)

    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        table_text = markdown_table_from_rows(rows)
        if table_text:
            parts.append("Word表格：\n" + table_text)

    return clean_text("\n".join(parts))


def convert_doc_to_docx(file_bytes: bytes, filename: str) -> bytes:
    """把旧版 .doc 转为 .docx；优先用 LibreOffice，其次用 Windows Word。"""
    with TemporaryDirectory(ignore_cleanup_errors=True) as temp_dir:
        temp_path = Path(temp_dir)
        safe_name = re.sub(r"[^A-Za-z0-9_.\-\u4e00-\u9fff]", "_", filename)
        doc_path = temp_path / safe_name
        if doc_path.suffix.lower() != ".doc":
            doc_path = doc_path.with_suffix(".doc")
        doc_path.write_bytes(file_bytes)

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "docx", "--outdir", str(temp_path), str(doc_path)],
                capture_output=True,
                text=True,
                timeout=90,
            )
            converted_path = doc_path.with_suffix(".docx")
            if result.returncode == 0 and converted_path.exists():
                return converted_path.read_bytes()

        if os.name == "nt":
            docx_path = doc_path.with_suffix(".docx")
            script = (
                "$ErrorActionPreference='Stop';"
                "$word=New-Object -ComObject Word.Application;"
                "$word.Visible=$false;"
                "try {"
                f"$doc=$word.Documents.Open('{str(doc_path)}');"
                f"$doc.SaveAs2('{str(docx_path)}', 16);"
                "$doc.Close($false);"
                "} finally { $word.Quit() }"
            )
            result = subprocess.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
                capture_output=True,
                text=True,
                timeout=90,
            )
            if result.returncode == 0 and docx_path.exists():
                return docx_path.read_bytes()

    raise RuntimeError(
        "当前环境不能直接解析旧版 .doc 文件。请在本机安装 Microsoft Word 或 LibreOffice，"
        "或者先把文件另存为 .docx 后再上传。"
    )


def extract_doc(file_bytes: bytes, filename: str) -> str:
    """提取旧版 .doc 文档文本：先转换为 .docx，再复用 Word 提取流程。"""
    docx_bytes = convert_doc_to_docx(file_bytes, filename)
    return extract_docx(io.BytesIO(docx_bytes))


def extract_pptx(file_obj) -> str:
    """提取 PowerPoint 中的文本、表格和备注内容。"""
    presentation = Presentation(file_obj)
    parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                table_text = markdown_table_from_rows(rows)
                if table_text:
                    slide_parts.append("PPT表格：\n" + table_text)
        try:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text.strip():
                slide_parts.append(f"备注：{notes_text_frame.text.strip()}")
        except Exception:
            pass
        if slide_parts:
            parts.append(f"第{slide_index}页\n" + "\n".join(slide_parts))

    return clean_text("\n\n".join(parts))


def convert_ppt_to_pptx(file_bytes: bytes, filename: str) -> bytes:
    """把旧版 .ppt 转为 .pptx；优先用 LibreOffice，其次用 Windows PowerPoint。"""
    with TemporaryDirectory(ignore_cleanup_errors=True) as temp_dir:
        temp_path = Path(temp_dir)
        safe_name = re.sub(r"[^A-Za-z0-9_.\-\u4e00-\u9fff]", "_", filename)
        ppt_path = temp_path / safe_name
        if ppt_path.suffix.lower() != ".ppt":
            ppt_path = ppt_path.with_suffix(".ppt")
        ppt_path.write_bytes(file_bytes)

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(temp_path), str(ppt_path)],
                capture_output=True,
                text=True,
                timeout=90,
            )
            converted_path = ppt_path.with_suffix(".pptx")
            if result.returncode == 0 and converted_path.exists():
                return converted_path.read_bytes()

        if os.name == "nt":
            pptx_path = ppt_path.with_suffix(".pptx")
            script = (
                "$ErrorActionPreference='Stop';"
                "$powerpoint=New-Object -ComObject PowerPoint.Application;"
                "try {"
                f"$presentation=$powerpoint.Presentations.Open('{str(ppt_path)}', $true, $false, $false);"
                f"$presentation.SaveAs('{str(pptx_path)}', 24);"
                "$presentation.Close();"
                "} finally { $powerpoint.Quit() }"
            )
            result = subprocess.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", script],
                capture_output=True,
                text=True,
                timeout=90,
            )
            if result.returncode == 0 and pptx_path.exists():
                return pptx_path.read_bytes()

    raise RuntimeError(
        "当前环境不能直接解析旧版 .ppt 文件。请在本机安装 PowerPoint 或 LibreOffice，"
        "或者先把文件另存为 .pptx 后再上传。"
    )


def extract_ppt(file_bytes: bytes, filename: str) -> str:
    """提取旧版 .ppt 文档文本：先转换为 .pptx，再复用 PPTX 提取流程。"""
    pptx_bytes = convert_ppt_to_pptx(file_bytes, filename)
    return extract_pptx(io.BytesIO(pptx_bytes))


def extract_image(file_obj, ocr_mode: str = "standard") -> str:
    """读取图片文件并按所选 OCR 模式执行识别。"""
    image = Image.open(file_obj)
    ocr_text = clean_ocr_text(run_ocr(image, mode=ocr_mode))
    vision_text = ""
    if st.session_state.get("pdf_extract_mode") == "vision":
        try:
            vision_text = call_vision_image_extraction(image)
        except Exception:
            vision_text = ""
    if vision_text:
        return clean_text(f"图片结构化识别：\n{vision_text}\n\n图片OCR文本：\n{ocr_text}")
    return ocr_text


def extract_plain_text(file_bytes: bytes) -> str:
    """读取 txt、md、json、jsonl 等纯文本材料。"""
    for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            return clean_text(file_bytes.decode(encoding))
        except UnicodeDecodeError:
            continue
    return clean_text(file_bytes.decode("utf-8", errors="ignore"))


def extract_tabular(file_obj, suffix: str) -> str:
    """把 CSV/Excel 表格转换为可分析的文本摘要。"""
    if suffix == "csv":
        df = pd.read_csv(file_obj)
    else:
        df = pd.read_excel(file_obj)
    preview = markdown_table_from_rows([list(df.columns), *df.head(30).astype(str).values.tolist()])
    columns = "、".join(str(col) for col in df.columns)
    return "\n".join(
        [
            f"表格数据，共{len(df)}行、{len(df.columns)}列。",
            f"字段：{columns}",
            "前30行预览：",
            preview,
        ]
    )


def split_sentences(text: str) -> list[str]:
    """按中英文句末标点切分句子，并过滤过短片段。"""
    sentences = re.split(r"(?<=[。！？.!?])\s*", text)
    return [sentence.strip() for sentence in sentences if len(sentence.strip()) >= 12]


def sentence_signature(sentence: str) -> str:
    """生成句子的去标点指纹，用于去重。"""
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "", sentence).lower()


def truncate_text(text: str, limit: int) -> str:
    """按指定长度截断文本，并尽量停在自然标点处。"""
    text = text.strip()
    if len(text) <= limit:
        return text
    cut = text[:limit].rstrip("，,；;。.!?、 ")
    for separator in ["。", "；", ";", ".", "，", ",", " "]:
        index = cut.rfind(separator)
        if index >= max(30, int(limit * 0.65)):
            cut = cut[:index]
            break
    return cut.rstrip("，,；;。.!?、 ")


def ensure_sentence_end(text: str) -> str:
    """确保摘要片段以句末标点结束。"""
    text = text.strip()
    if not text:
        return text
    return text if text[-1] in "。.!?！？" else text + "。"


def normalize_summary_input(text: str) -> str:
    """去掉训练样本或提示词包装，只保留适合摘要的正文部分。"""
    text = clean_text(text)
    if "Abstract:" in text:
        return text[text.find("Abstract:"):]
    markers = ["核心内容：", "核心内容:"]
    for marker in markers:
        if marker in text:
            return text[text.find(marker) + len(marker):].strip()
    return text


def tokenize(text: str) -> list[str]:
    """使用 jieba 分词并过滤停用词、纯数字和符号。"""
    words = []
    for word in jieba.cut(text):
        word = word.strip().lower()
        if len(word) < 2:
            continue
        if word in STOPWORDS:
            continue
        if re.fullmatch(r"[\d\W_]+", word):
            continue
        words.append(word)
    return words


def select_key_sentences(text: str, max_sentences: int = 5) -> list[str]:
    """基于词频、位置和长度评分，从原文中抽取关键句。"""
    raw_sentences = split_sentences(text)
    sentences = []
    seen = set()
    for sentence in raw_sentences:
        signature = sentence_signature(sentence)
        if not signature or signature in seen:
            continue
        seen.add(signature)
        sentences.append(sentence)
    if not sentences:
        return [text[:500]] if text.strip() else []

    word_freq = Counter(tokenize(text))
    scored = []
    for index, sentence in enumerate(sentences):
        score = sum(min(word_freq.get(word, 0), 8) for word in tokenize(sentence))
        if re.search(r"[\u4e00-\u9fff]", sentence):
            score *= 1.15
        position_bonus = 1.15 if index < max(3, len(sentences) // 8) else 1.0
        length_penalty = 1.0 if 35 <= len(sentence) <= 180 else 0.75
        score = score / max(len(sentence), 1) * position_bonus * length_penalty
        scored.append((score, index, sentence))

    selected = sorted(scored, reverse=True)[:max_sentences]
    selected = sorted(selected, key=lambda item: item[1])
    return [item[2] for item in selected]


def compact_sentences(sentences: list[str], limit: int = 220) -> str:
    """把若干关键句压缩成一段较自然的本地摘要片段。"""
    unique_sentences = []
    seen = set()
    for sentence in sentences:
        signature = sentence_signature(sentence)
        if not signature or signature in seen:
            continue
        seen.add(signature)
        unique_sentences.append(sentence)
    content = "；".join(sentence.rstrip("。；;") for sentence in unique_sentences if sentence.strip())
    content = re.sub(r"\s+", " ", content).strip("；; ")
    if not content:
        return "原文有效信息较少，暂无法生成完整摘要。"
    if len(content) > limit:
        content = truncate_text(content, limit)
    return content + "。"


ENGLISH_TERM_MAP = {
    "document": "文档",
    "intelligence": "智能处理",
    "ocr": "OCR识别",
    "pdf": "PDF解析",
    "word": "Word解析",
    "summarization": "摘要生成",
    "summary": "摘要",
    "cloud": "云端部署",
    "deployment": "部署",
    "streamlit": "Streamlit系统",
    "tesseract": "Tesseract OCR",
    "user": "用户隔离",
    "sqlite": "SQLite存储",
    "model": "模型",
    "experiment": "实验",
    "results": "结果验证",
    "contract": "合同",
    "payment": "付款金额",
    "service": "服务期限",
    "breach": "违约条款",
}


def local_keywords(text: str, limit: int = 8) -> str:
    """提取本地摘要使用的关键词，并把常见英文术语映射成中文表达。"""
    word_freq = Counter(tokenize(text))
    ranked_words = [word for word, _ in word_freq.most_common(40)]
    chinese_words = [word for word in ranked_words if re.search(r"[\u4e00-\u9fff]", word)]
    mapped_words = []
    for word in ranked_words:
        mapped = ENGLISH_TERM_MAP.get(word.lower())
        if mapped and mapped not in mapped_words:
            mapped_words.append(mapped)
    selected = chinese_words[:limit] if chinese_words else mapped_words[:limit]
    if not selected:
        selected = ranked_words[:limit]
    return "、".join(selected) or "暂无明显关键词"


def find_section_text(text: str, names: list[str], limit: int = 260) -> str:
    """按章节标题提取论文常见区块，如摘要、方法、结果和结论。"""
    joined = clean_text(text)
    for name in names:
        pattern = rf"(?:^|\n)\s*{re.escape(name)}\s*[:：]?\s*(.*?)(?=\n\s*(?:摘要|abstract|关键词|keywords|引言|introduction|方法|methodology|method|实验|experiment|结果|results|讨论|discussion|结论|conclusion|references|参考文献|[一二三四五六七八九十]、|\d+[.、])|$)"
        match = re.search(pattern, joined, flags=re.IGNORECASE | re.DOTALL)
        if match:
            content = clean_text(match.group(1))
            if len(content) >= 20:
                return compact_sentences(split_sentences(content)[:3] or [content[:limit]], limit=limit)
    return ""


def extract_sentences_by_terms(text: str, terms: list[str], limit: int = 220) -> str:
    """从文本中抽取包含指定风险词或行动词的句子。"""
    sentences = split_sentences(text)
    matched = [sentence for sentence in sentences if any(term.lower() in sentence.lower() for term in terms)]
    return compact_sentences(matched[:3], limit=limit) if matched else ""


def local_academic_summary(text: str, key_sentences: list[str], keywords: str) -> str:
    """在大模型不可用时，为论文类文档生成本地学术摘要。"""
    text = normalize_summary_input(text)
    key_sentences = select_key_sentences(text, max_sentences=7)
    abstract = find_section_text(text, ["摘要", "Abstract"], limit=260)
    method = find_section_text(text, ["方法", "Methodology", "Method"], limit=220)
    result = find_section_text(text, ["结果", "Results", "实验", "Experiment"], limit=220)
    conclusion = find_section_text(text, ["结论", "Conclusion", "Conclusions"], limit=220)

    if abstract:
        parts = [f"该材料主要讨论{keywords}等内容。{abstract}"]
        if method:
            parts.append(f"在分析路径上，{method}")
        if result:
            parts.append(f"主要结论是，{result}")
        if conclusion:
            parts.append(f"进一步来看，{conclusion}")
        return clean_summary_output(ensure_sentence_end(truncate_text("".join(parts), 650)))

    academic_body = compact_sentences(key_sentences[:5], limit=460)
    return clean_summary_output(f"研究聚焦于{keywords}等问题。{academic_body}")


def local_risk_or_action(text: str, terms: list[str], fallback: str) -> str:
    """提取风险或行动建议；没有命中时返回固定兜底说明。"""
    extracted = extract_sentences_by_terms(text, terms, limit=200)
    return extracted if extracted else fallback


def infer_material_style(filename: str = "", file_type: str = "") -> str:
    """根据文件后缀或路由结果判断材料形态，用于选择更合适的摘要模板。"""
    suffix = Path(filename).suffix.lower()
    if file_type == "presentation" or suffix in {".ppt", ".pptx"}:
        return "presentation"
    if file_type == "image" or suffix in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}:
        return "visual"
    return "document"


def local_presentation_summary(text: str, keywords: str, key_sentences: list[str], material_style: str) -> str:
    """为 PPT 或图片 OCR 结果生成更自然的本地兜底摘要。"""
    label = "PPT演示材料" if material_style == "presentation" else "图片/视觉材料"
    theme = compact_sentences(key_sentences[:2], limit=180)
    structure = compact_sentences(key_sentences[2:5], limit=240) if len(key_sentences) > 2 else "材料以标题、短句或图中文字为主，结构信息较为碎片化。"
    details = compact_sentences(key_sentences[5:8], limit=220) if len(key_sentences) > 5 else "未识别到足够连续的正文段落，摘要主要依据可读标题和关键短句生成。"
    return "\n".join(
        [
            f"材料主题：该{label}主要涉及{keywords}等内容。{theme}",
            f"内容结构：材料呈现了若干关键模块或步骤，核心线索包括{structure}",
            f"关键设定：可识别信息显示，材料重点强调{details}",
            "阅读提示：由于 PPT 或图片通常只包含标题、要点和图中文字，系统不会据此扩展实验效果、适用边界或行动建议；如需更完整摘要，建议同时上传配套讲稿、论文正文或说明文档。",
        ]
    )


def summarize_text(text: str, max_sentences: int = 5, summary_mode: str = "中文摘要", material_style: str = "document") -> str:
    """本地规则摘要入口，用于无 API 或大模型失败时兜底。"""
    text = normalize_summary_input(text)
    key_sentences = select_key_sentences(text, max_sentences=max_sentences + 4)
    if not key_sentences:
        return "原文有效信息较少，暂无法生成完整摘要。"

    category, _ = classify_document(text)
    keywords = local_keywords(text)
    if material_style in {"presentation", "visual"}:
        summary = local_presentation_summary(text, keywords, key_sentences, material_style)
        if summary_mode != "中文摘要":
            return f"本地规则摘要暂不支持稳定的{summary_mode}生成，以下为中文兜底摘要：\n{summary}"
        return summary

    overview = compact_sentences(key_sentences[:2], limit=180)
    details_sentences = key_sentences[2:5]
    conclusion_sentences = key_sentences[5:7]
    details = compact_sentences(details_sentences, limit=220) if details_sentences else "原文未提供更多可区分的关键信息。"
    conclusion = compact_sentences(conclusion_sentences, limit=180) if conclusion_sentences else "原文未提供独立于前文的明确结论。"

    if is_academic_document(text):
        summary = local_academic_summary(text, key_sentences, keywords)
        if summary_mode != "中文摘要":
            return f"本地规则摘要暂不支持稳定的{summary_mode}生成，以下为中文兜底摘要：\n{summary}"
        return summary

    risks = local_risk_or_action(
        text,
        ["风险", "注意", "限制", "问题", "违约", "失败", "缺少", "不足", "risk", "limitation", "breach", "missing"],
        "原文未提出明确风险或注意事项。",
    )
    actions = local_risk_or_action(
        text,
        ["建议", "应当", "需要", "必须", "请", "补充", "确认", "部署", "安装", "should", "need", "must", "install", "deploy"],
        "原文未提出明确建议行动。",
    )

    summary = "\n".join(
        [
            f"核心概述：该材料属于{category}，重点涉及{keywords}等内容。{overview}",
            f"关键信息：{details}",
            f"重要结论：{conclusion}",
            f"风险或注意事项：{risks}",
            f"建议行动：{actions}",
        ]
    )
    if summary_mode != "中文摘要":
        return f"本地规则摘要暂不支持稳定的{summary_mode}生成，以下为中文兜底摘要：\n{summary}"
    return summary


def get_secret(name: str) -> str | None:
    """优先从 Streamlit secrets 读取配置，读取不到再使用环境变量。"""
    try:
        value = st.secrets.get(name)
    except Exception:
        value = None
    return value or os.getenv(name)


def get_llm_config() -> tuple[str | None, str, list[str]]:
    """读取大模型 API 配置，并组装主模型和备用模型列表。"""
    api_key = get_secret("LLM_API_KEY")
    api_base = get_secret("LLM_API_BASE") or "https://api.openai.com/v1"
    model = get_secret("LLM_MODEL") or "gpt-5.4-mini"
    backup_models = get_secret("LLM_BACKUP_MODELS") or ""
    models = [model]
    for item in re.split(r"[,;\n]+", backup_models):
        item = item.strip()
        if item and item not in models:
            models.append(item)
    return api_key, api_base.rstrip("/"), models


def get_vision_llm_config() -> tuple[str | None, str, str]:
    """读取 PDF 多模态提取使用的视觉模型配置；未单独配置时复用摘要模型。"""
    api_key = get_secret("VISION_API_KEY") or get_secret("LLM_API_KEY")
    api_base = get_secret("VISION_API_BASE") or get_secret("LLM_API_BASE") or "https://api.openai.com/v1"
    model = get_secret("VISION_MODEL") or get_secret("LLM_MODEL") or "gpt-5.4-mini"
    return api_key, api_base.rstrip("/"), model


@st.cache_resource(show_spinner=False)
def load_local_generation_pipeline(model_path: str):
    """按需加载本地生成/翻译模型；未安装依赖或无模型时返回错误信息。"""
    try:
        from transformers import pipeline
    except Exception as exc:
        return None, f"未安装本地模型依赖 transformers/torch：{exc}"
    path = Path(model_path)
    if not path.exists():
        return None, f"本地模型路径不存在：{model_path}"
    try:
        generator = pipeline("text2text-generation", model=str(path), tokenizer=str(path))
        return generator, ""
    except Exception as exc:
        return None, f"本地模型加载失败：{exc}"


def get_local_model_path() -> str | None:
    """读取可选本地英文/双语摘要模型路径。"""
    return get_secret("LOCAL_SUMMARY_MODEL_PATH") or os.getenv("LOCAL_SUMMARY_MODEL_PATH")


def local_model_summary(text: str, summary_mode: str, academic: bool) -> tuple[str | None, str]:
    """使用本地模型生成英文或双语摘要；没有模型时返回原因。"""
    model_path = get_local_model_path()
    if not model_path:
        return None, "未配置 LOCAL_SUMMARY_MODEL_PATH。"
    generator, error = load_local_generation_pipeline(model_path)
    if generator is None:
        return None, error
    if summary_mode == "英文摘要":
        instruction = "Write an English literature-review style summary." if academic else "Write an English business summary."
    elif summary_mode == "双语摘要":
        instruction = "Write a bilingual Chinese and English summary."
    else:
        instruction = "写一段中文摘要。"
    prompt = f"{instruction}\n\nSource:\n{build_llm_source(text, academic, limit=3500)}"
    try:
        result = generator(prompt, max_new_tokens=420, do_sample=False)
        generated = result[0].get("generated_text") or result[0].get("summary_text") or ""
        generated = clean_summary_output(generated)
        if generated:
            return generated, ""
        return None, "本地模型返回为空。"
    except Exception as exc:
        return None, f"本地模型生成失败：{exc}"



def parse_llm_response(response: requests.Response) -> str:
    """兼容普通 JSON 和流式 SSE 两种 OpenAI 风格响应。"""
    content_type = response.headers.get("content-type", "")
    if "text/event-stream" not in content_type:
        data = response.json()
        return data["choices"][0]["message"]["content"].strip()

    chunks = []
    for raw_line in response.iter_lines(decode_unicode=True):
        if not raw_line:
            continue
        line = raw_line.strip()
        if not line.startswith("data:"):
            continue
        payload = line.removeprefix("data:").strip()
        if payload == "[DONE]":
            break
        try:
            data = json.loads(payload)
        except json.JSONDecodeError:
            continue
        choices = data.get("choices") or []
        if not choices:
            continue
        choice = choices[0]
        delta = choice.get("delta") or {}
        message = choice.get("message") or {}
        chunks.append(delta.get("content") or message.get("content") or "")
    return "".join(chunks).strip()



def call_llm_summary(api_key: str, api_base: str, models: str | list[str], prompt: str, max_tokens: int = 900, timeout: int = 45) -> tuple[str, str]:
    """调用大模型生成摘要；支持自动重试和备用模型切换。"""
    model_list = [models] if isinstance(models, str) else models
    last_error = None
    for model in model_list:
        for attempt in range(1, 3):
            try:
                response = requests.post(
                    f"{api_base}/chat/completions",
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": model,
                        "messages": [
                            {
                                "role": "system",
                                "content": "你是一个严谨的中文文档分析助手。只依据原文总结，不编造，不把乱码当正文。",
                            },
                            {"role": "user", "content": prompt},
                        ],
                        "temperature": 0.1,
                        "max_tokens": max_tokens,
                        "stream": True,
                    },
                    stream=True,
                    timeout=timeout,
                )
                response.raise_for_status()
                content = parse_llm_response(response)
                if content.strip():
                    return content, model
                last_error = RuntimeError("empty model response")
            except Exception as exc:
                last_error = exc
                if attempt < 2:
                    time.sleep(1.5 * attempt)
    raise RuntimeError(f"LLM request failed after retries/models: {last_error}")


def is_academic_document(text: str) -> bool:
    """根据章节标题和论文信号判断文档是否属于学术/论文类。"""
    lower = text.lower()
    if "abstract:" in lower or "abstract\n" in lower or "摘要：" in text or "摘要:" in text:
        return True
    if ("introduction:" in lower or "conclusion:" in lower) and any(
        term in lower for term in ["paper", "study", "method", "model", "experiment", "result", "transformer"]
    ):
        return True
    signals = [
        "abstract", "keywords", "introduction", "method", "methodology",
        "experiment", "results", "discussion", "conclusion", "references",
        "摘要", "关键词", "引言", "方法", "实验",
        "结果", "讨论", "结论", "参考文献",
    ]
    return sum(1 for signal in signals if signal in lower) >= 3


SUMMARY_POLLUTION_PHRASES = [
    "作为ai",
    "as an ai",
    "如果你需要",
    "我可以继续",
    "下一步",
    "你可以直接",
    "请提供",
    "重新发",
    "无法可靠",
    "无法还原",
    "我看到你",
    "输出语言",
    "输出要求",
    "审稿目标",
    "硬约束",
    "json字段",
    "当前摘要",
    "以下是",
    "以下为",
    "根据你的要求",
    "根据原文要求",
    "本地规则摘要暂不支持",
    "未调用在线质量复核",
    "未完成在线pdca",
    "final review",
    "pdca",
    "revision_guidance",
    "chinese_final",
    "english_final",
    "```",
]


def has_summary_pollution(summary: str) -> bool:
    """识别不应出现在最终摘要中的提示词、JSON、审稿过程和系统说明。"""
    text = summary.strip()
    if not text:
        return False
    lower = text.lower()
    # 第一层用短语拦截明显的提示词残留和系统过程说明。
    if any(phrase in lower for phrase in SUMMARY_POLLUTION_PHRASES):
        return True
    # 第二层只拦截 "score:" 这类 JSON 字段，避免误伤正常英文里的 risk score。
    json_key_pattern = r'["“]?(?:score|issues|revision_guidance|chinese_final|english_final|verdict|level)["”]?\s*[:：]'
    if re.search(json_key_pattern, text, flags=re.IGNORECASE):
        return True
    # 第三层按行首判断审稿报告残留，防止 Final Review 内容混入摘要正文。
    process_markers = [
        r"^审稿",
        r"^评分[:：]",
        r"^主要问题[:：]",
        r"^修订建议[:：]",
        r"^输出[:：]",
        r"^原文[:：]",
        r"^材料类型[:：]",
        r"^任务上下文[:：]",
    ]
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return any(any(re.search(pattern, line, flags=re.IGNORECASE) for pattern in process_markers) for line in lines)


def strip_summary_pollution(summary: str) -> str:
    """删除明显不属于摘要正文的过程说明、JSON字段和提示词残留。"""
    text = clean_text(summary)
    if not text:
        return text

    # 如果模型误把 JSON 直接当摘要返回，优先提取可用的中英文终稿字段。
    json_match = re.search(r"\{.*\}", text, flags=re.DOTALL)
    if json_match:
        try:
            data = json.loads(json_match.group(0))
            chinese = clean_text(str(data.get("chinese_final", "")))
            english = clean_text(str(data.get("english_final", "")))
            if chinese and english:
                return f"{chinese}\n\nEnglish Summary\n{english}"
            if chinese:
                return chinese
            if english:
                return english
        except Exception:
            pass

    polluted_line_patterns = [
        r"^```",
        r"^以下是",
        r"^以下为",
        r"^根据你的要求",
        r"^输出要求",
        r"^审稿目标",
        r"^硬约束",
        r"^JSON字段",
        r"^当前摘要[:：]",
        r"^原文[:：]",
        r"^材料类型[:：]",
        r"^任务上下文[:：]",
        r"^Final Review[:：]",
        r"^PDCA",
        r"^评分[:：]",
        r"^主要问题[:：]",
        r"^修订建议[:：]",
        r'^["“]?(?:score|issues|revision_guidance|chinese_final|english_final|verdict|level)["”]?\s*[:：]',
    ]
    cleaned_lines = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        # 只删除整行污染内容，不在句子内部做激进删除，减少误删真实摘要信息。
        if any(re.search(pattern, line, flags=re.IGNORECASE) for pattern in polluted_line_patterns):
            continue
        cleaned_lines.append(line)
    text = "\n".join(cleaned_lines) if cleaned_lines else text
    text = re.sub(r"^中文摘要\s*[:：]\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^摘要\s*[:：]\s*", "", text, flags=re.IGNORECASE)
    return text.strip()


def is_low_quality_summary(summary: str, academic: bool = False, summary_mode: str = "中文摘要") -> bool:
    """检查摘要是否过短、重复、对话化、乱码化或缺少学术摘要要素。"""
    lines = [line.strip() for line in summary.splitlines() if line.strip()]
    if not lines:
        return True
    if has_summary_pollution(summary):
        return True
    normalized = [re.sub(r"\W+", "", line.lower()) for line in lines]
    normalized = [line for line in normalized if line]
    if len(normalized) >= 3 and len(set(normalized)) <= max(1, len(normalized) // 3):
        return True
    if len(summary.strip()) < 120:
        return True
    bad_phrases = [
        "研究主题：",
        "需要强调",
        "原文说明",
        "原文指出",
        "原文提到",
        "原文表示",
        "原文介绍",
        "该文档提到",
        "该文档说明",
        "本文旨在",
        "本文主要围绕",
        "这段内容主要讨论",
        "this document",
        "the source",
    ]
    if any(phrase.lower() in summary.lower() for phrase in bad_phrases):
        return True
    if summary_mode != "中文摘要":
        # 英文/双语模式先检查英文体量，防止模型只输出中文兜底摘要。
        latin_words = re.findall(r"[A-Za-z]{3,}", summary)
        if summary_mode == "英文摘要":
            return len(latin_words) < 35
        if summary_mode == "双语摘要":
            has_chinese = bool(re.search(r"[\u4e00-\u9fff]", summary))
            has_english = len(latin_words) >= 25
            return not (has_chinese and has_english)
    if academic:
        # 学术摘要至少要覆盖问题、方法、结果/证据、贡献中的多数要素。
        requirement_groups = [
            ["研究", "问题", "场景", "背景", "缺口", "需求", "面向"],
            ["方法", "模型", "系统", "框架", "采用", "构建", "设计"],
            ["结果", "表明", "发现", "验证", "测试", "样本外", "Precision", "Recall", "PR-AUC", "Lift"],
            ["贡献", "边际贡献", "创新", "价值", "含义", "启示"],
        ]
        if sum(any(word in summary for word in group) for group in requirement_groups) < 3:
            return True
    return False


def needs_top_venue_polish(summary: str, source_text: str = "") -> bool:
    """判断学术摘要是否还停留在普通概述，而不是顶刊/顶会式研究摘要。"""
    text = summary.strip()
    if not text:
        return False
    lower_text = text.lower()
    template_markers = [
        "核心概述",
        "关键信息",
        "重要结论",
        "建议行动",
        "该材料属于",
        "该文档",
        "这份材料",
        "主要讨论",
        "主要围绕",
        "具有参考意义",
        "一定参考价值",
    ]
    if any(marker in text for marker in template_markers):
        return True

    structure_groups = [
        ["背景", "缺口", "需求", "问题", "有限审查", "类别不平衡", "capital market", "research gap"],
        ["提出", "构建", "设计", "重构", "框架", "方法", "模型", "propose", "framework", "method"],
        ["样本", "数据", "实证", "实验", "测试", "时间阻断", "PR-AUC", "Precision", "Recall", "Lift", "evaluation", "out-of-sample"],
        ["结果", "表明", "显示", "达到", "提升", "发现", "results", "show", "achieve"],
        ["贡献", "边际贡献", "创新", "启示", "含义", "contribution", "implication"],
    ]
    hits = sum(any(term.lower() in lower_text for term in group) for group in structure_groups)
    source_has_results = any(term.lower() in source_text.lower() for term in ["pr-auc", "precision", "recall", "lift", "accuracy", "bacc", "结果", "测试集"])
    if source_has_results and hits < 4:
        return True
    return hits < 3


def needs_style_polish(summary: str) -> bool:
    """判断摘要是否存在句子过长、模板味太重或符号残留等风格问题。"""
    text = summary.strip()
    if not text:
        return False
    if has_summary_pollution(text):
        return True
    sentences = [sentence for sentence in re.split(r"[。！？!?]\s*", text) if sentence.strip()]
    if any(len(sentence) > 150 for sentence in sentences):
        return True
    dense_markers = ["原文说明", "原文指出", "该文档提到", "本文旨在", "围绕", "系统梳理", "进一步分析", "归纳了", "协同作用", "相关内容为", "理论参考"]
    if sum(text.count(marker) for marker in dense_markers) >= 3:
        return True
    if text.count("、") >= 12 and len(sentences) <= 4:
        return True
    if any(symbol in text for symbol in ["→", "=>", "＋", " + "]):
        return True
    return False


def clean_summary_output(summary: str) -> str:
    """对最终摘要做格式清洗，去除箭头、加号、括号和异常标点。"""
    text = strip_summary_pollution(summary)
    if not text:
        return text
    replacements = {
        "→": "，",
        "⇒": "，",
        "➜": "，",
        "＋": "和",
        "+": "和",
        " + ": "和",
        "；解释：": "。具体而言，",
        "解释：": "具体而言，",
        "对应的两部分内容": "其内容可以分为两个层面",
        "原文说明": "",
        "原文指出": "",
        "原文提到": "",
        "原文表示": "",
        "原文介绍": "",
        "原文阐述": "",
        "该文档提到": "",
        "该文档说明": "",
        "该材料提到": "",
        "该材料说明": "",
        "这段内容主要讨论": "研究聚焦于",
        "本文旨在研究": "研究聚焦于",
        "本文旨在": "研究聚焦于",
        "本文主要围绕": "研究聚焦于",
        "本文围绕": "研究聚焦于",
        "主要围绕": "重点涉及",
        "围绕": "聚焦于",
        "旨在研究": "聚焦于",
        "研究聚焦于研究": "研究聚焦于",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    text = re.sub(r"^(?:，|。|：|:|；|;)+", "", text)
    text = re.sub(r"(?:^|。)(?:该文档|这份文档|该材料|这份材料)(?:主要)?(?:围绕|讨论|介绍|阐述)", "。内容聚焦于", text)
    text = re.sub(r"(?:^|。)(?:文章|文本)(?:主要)?(?:围绕|讨论|介绍|阐述)", "。内容聚焦于", text)
    text = re.sub(r"^。", "", text)
    text = re.sub(r"[」』》]+", "", text)
    text = re.sub(r"[「『《]+", "", text)
    text = re.sub(r"[（）()【】\[\]]", "", text)
    text = re.sub(r"\s*([，。！？；：、])\s*", r"\1", text)
    text = re.sub(r"，{2,}", "，", text)
    text = re.sub(r"。{2,}", "。", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()


def clean_english_summary_output(summary: str) -> str:
    """清洗英文摘要，保留英文标点，避免出现中文句号、分号和逗号。"""
    text = strip_summary_pollution(summary)
    if not text:
        return text
    replacements = {
        "；": "; ",
        "，": ", ",
        "。": ". ",
        "：": ": ",
        "（": "(",
        "）": ")",
        "＋": " and ",
        "→": ", ",
        "⇒": ", ",
        "➜": ", ",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    text = re.sub(r"([,.;:!?])(?=[A-Za-z])", r"\1 ", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()


def build_summary_prompt(source_text: str, academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """根据文档类型构建大模型摘要提示词。"""
    output_spec = summary_output_spec(summary_mode, academic, material_style, cn_target, en_target)
    material_label = prompt_material_label(academic, material_style)
    language_rule = prompt_language_contract(summary_mode)
    evidence_rules = prompt_evidence_contract(material_style)
    if material_style in {"presentation", "visual"}:
        return f"""你是一名严谨的材料解读助手。请把提取文本整理成适合阅读的{material_label}摘要。

任务上下文：
- 材料类型：{material_label}
- 输出语言：{language_rule}
- 输出合同：{output_spec}

生成前请在内部完成判断，但不要输出判断过程：
- 识别材料主题、页面/图示结构、关键机制、变量、公式、流程或指标。
- 区分“可确认事实”和“无法从材料确认的推断”。
- 对碎片化标题和图中文字进行合并，不把它们扩展成原文没有的结论。

证据与风格约束：
{evidence_rules}

提取文本：
{source_text}
"""
    if academic:
        top_venue_contract = top_venue_abstract_contract(summary_mode)
        return f"""你是一名熟悉顶级期刊和顶级会议写作规范的学术摘要编辑。请根据原文生成一段可用于论文投稿或项目展示的高质量研究摘要。

任务上下文：
- 材料类型：{material_label}
- 输出语言：{language_rule}
- 输出合同：{output_spec}

顶刊/顶会摘要风格合同：
{top_venue_contract}

生成前请在内部完成判断，但不要输出判断过程：
- 提取研究背景/问题、方法或系统设计、数据/实验/应用场景、主要发现、创新点或贡献、应用意义。
- 判断哪些结论有原文支持，哪些只是作者设想或方法描述。
- 原文缺少定量结果时，只写定性发现；原文信息不足时，保持概括，不补细节。
- 如果原文是方法设计或论文草稿，要把摘要写成“问题驱动的研究贡献”，不要写成文档说明。
- 如果原文包含模型指标、Top-K、PR-AUC、时间阻断、样本外检验等证据，请优先保留这些能支撑贡献的关键信息。

证据与风格约束：
{evidence_rules}
- 用正式、凝练、连贯的学术表达，不要写成“本文属于某类文档”的分类说明。
- 每句话只表达一个主要意思，句子之间要有因果、递进或转折关系。
- 避免“围绕、系统梳理、进一步分析、归纳了、研究表明、相关内容为”等模板化表达。
- 不要为了模仿顶刊而夸大结论；如果原文只说明设定、框架或实验设计，就把贡献限定在任务重构、机制设计和评估方案。

以下顶刊/顶会风格范例只参考组织方式、信息密度和学术语气，不参考事实：
{format_academic_examples(source_text)}

原文：
{source_text}
"""

    return f"""你是一名专业文档分析员，请根据原文生成最终摘要。

任务上下文：
- 材料类型：{material_label}
- 输出语言：{language_rule}
- 输出合同：{output_spec}

生成前请在内部完成判断，但不要输出判断过程：
- 提取主体、事项、日期、金额、要求、结论、风险、限制和行动项。
- 判断哪些是明确事实，哪些只是背景描述或重复信息。
- 没有风险或行动项时，按输出合同简洁说明，不额外创造建议。

证据与风格约束：
{evidence_rules}
- 每个小标题下写 1-2 句完整句子。
- 语言要像人工业务摘要，优先写清楚“材料在说什么、为什么重要、最后形成什么判断”。
- 合并重复信息，避免空泛评价和名词堆叠。

原文：
{source_text}
"""


def build_pdca_plan(academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """生成 PDCA 摘要质量检查中的 Plan 标准。"""
    output_spec = summary_output_spec(summary_mode, academic, material_style, cn_target, en_target)
    if material_style in {"presentation", "visual"}:
        return (
            f"Plan质量标准：{output_spec} 摘要应围绕材料主题、内容结构、关键要点和阅读提示组织；"
            "只依据可读文本，不把碎片化标题扩展成原文没有的实验结论、风险或行动建议；关键判断必须能在原文中找到直接依据。"
        )
    if academic:
        return (
            f"Plan质量标准：{output_spec} 按顶刊/顶会摘要结构组织，覆盖背景缺口或现实问题、"
            "研究问题、方法或系统设计、数据/实验/应用场景、主要支持性发现、边际贡献和应用含义；"
            "只依据原文，避免幻觉、夸大、绝对化表述和技术贡献/实验结果混淆；每个关键结论都必须有原文直接支撑。"
        )
    return (
        f"Plan质量标准：{output_spec} 覆盖核心概述、关键信息、"
        "重要结论、风险或注意事项、建议行动；只依据原文，保留关键事实，"
        "避免新增原文没有的判断、风险和行动项；无法在原文中定位依据的内容一律删除或降级表述。"
    )


def pdca_check_summary_prompt(source_text: str, summary: str, academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """构建 PDCA 的 Check 提示词，让大模型审查摘要是否忠实可靠。"""
    expected = "演示/视觉材料摘要" if material_style in {"presentation", "visual"} else "中文学术摘要" if academic else "五段式中文业务摘要"
    material_label = prompt_material_label(academic, material_style)
    return f"""你是严格的摘要质量审查员。请执行PDCA循环中的Check环节，只依据原文审查当前摘要。

Plan质量标准：
{build_pdca_plan(academic, summary_mode, material_style, cn_target, en_target)}

审查对象：
- 材料类型：{material_label}
- 期望摘要类型：{expected}

审查维度：
1. 原文没有支持的结论、数据、模型、指标、年份或政策建议。
2. 过度夸大的评价，如“证明”“显著领先”“企业级部署”“全面解决”等缺少充分依据的表述。
3. 模糊但绝对化的表述，如“最佳”“最优”“完全”“必然”“长期有效”。
4. 遗漏关键实验条件、数据集、评价指标、适用范围或限制。
5. 混淆技术贡献与实验结果，或把作者主张写成已被充分验证的事实。
6. 摘要结构、长度和表达是否符合{expected}要求。
7. 如果材料是PPT或图片，检查摘要是否强行补写了原文没有的实验结论、风险或行动建议。
8. 如果输出要求是双语，检查中文和英文是否分段清楚，是否出现中英文混在同一摘要框的风险。
9. 检查摘要中是否混入JSON字段、审稿意见、评分、提示词残留、系统处理说明、对用户的操作建议或“以下是/根据你的要求”等非摘要正文。
10. 逐句核对摘要：每个关键句必须能在原文中找到直接支持；找不到依据的句子必须列入issues。
11. 严格核对数字、样本量、模型名、指标、比较对象、因果关系和贡献声明；原文没有明确出现时判为幻觉或过度推断。
12. 不允许用“可能表明、具有参考意义、重要启示、建议行动”等泛化话术包装原文没有的判断。

评分依据：
- 忠实性40分：摘要中的事实、结论、数据和因果关系必须能被原文支持，不能幻觉或夸大。
- 完整性25分：覆盖核心研究问题、数据/方法、主要发现、限制或启示，不能遗漏关键条件。
- 结构与任务匹配20分：符合指定语言、长度、段落结构和摘要类型。
- 表达质量15分：语言自然、凝练、专业，避免机械堆词、重复和对话式表达。

输出要求：
- 只输出JSON，不要Markdown。
- JSON字段包括：score、verdict、issues、revision_guidance。
- score为0-100整数；verdict只能是pass或revise。
- 如果score低于{pass_score}，verdict必须是revise。
- 只要发现无依据关键句、虚构数字、虚构模型/指标或混入非摘要正文，verdict必须是revise，score不得高于{max(pass_score - 1, 0)}。
- issues是字符串数组；没有问题时输出空数组。
- 有问题时issues列出3-6个最重要问题，明确指出哪些内容需要删除、降级或改写。
- revision_guidance用一段中文说明如何修改，必须具体说明删掉什么、保留什么，以及哪些内容应改为“原文未说明”或更保守表述。

原文：
{source_text[:9000]}

当前摘要：
{summary}
"""


def parse_pdca_review(review_text: str, pass_score: int = 88) -> dict:
    """解析 PDCA Check 阶段返回的 JSON 审查结果。"""
    fallback = {
        "score": 0,
        "verdict": "revise",
        "issues": ["审查结果解析失败，已按需修订处理。"],
        "revision_guidance": clean_text(review_text)[:500] or "请按原文重新检查并修订摘要。",
    }
    if not review_text:
        return fallback
    cleaned = review_text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, flags=re.DOTALL)
    if match:
        cleaned = match.group(0)
    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        return fallback
    try:
        score = int(data.get("score", 0))
    except (TypeError, ValueError):
        score = 0
    verdict = str(data.get("verdict", "revise")).strip().lower()
    if verdict not in {"pass", "revise"}:
        verdict = "pass" if score >= pass_score else "revise"
    issues = data.get("issues") or []
    if isinstance(issues, str):
        issues = [issues]
    issues = [clean_text(str(item)) for item in issues if clean_text(str(item))]
    return {
        "score": max(0, min(100, score)),
        "verdict": verdict,
        "issues": issues,
        "revision_guidance": clean_text(str(data.get("revision_guidance", ""))),
    }


def top_venue_final_review_prompt(source_text: str, summary: str, summary_mode: str, cn_target: int | None = None, en_target: int | None = None) -> str:
    """构建最终顶刊/顶会摘要审稿和中英文改写提示词。"""
    return f"""你是顶级期刊/顶级会议摘要审稿人与语言编辑。请模拟人工审稿流程，检查当前摘要并给出最终中文和英文摘要。

审稿目标：
- 像审稿人一样指出摘要是否清楚呈现背景缺口、研究任务、方法机制、数据/实验设置、主要结果、边际贡献和应用含义。
- 像语言编辑一样把摘要改成顶刊/顶会风格，要求克制、凝练、问题驱动、证据承接清楚。
- 中文和英文都要自动生成；英文不是逐字翻译，而是符合英文论文摘要习惯的改写。

硬约束：
- 只能依据原文和当前摘要，不得新增原文没有的数据、指标、样本、年份、结论或建议。
- 如果原文没有定量结果，只能写方法设计、评估方案或定性发现。
- 不要输出Markdown，不要输出解释过程，只输出JSON。
- JSON字段必须包括：score、level、issues、revision_guidance、chinese_final、english_final。
- score为0-100整数；level只能是excellent、usable、revise。
- issues为字符串数组；revision_guidance为中文短段落。
- chinese_final为300-500个汉字的中文顶刊/顶会风格摘要。
- english_final为160-240词的英文顶刊/顶会风格摘要。

原文：
{source_text[:9000]}

当前摘要：
{summary}
"""


def parse_top_venue_final_review(review_text: str) -> dict:
    """解析最终顶刊/顶会审稿与中英文改写结果。"""
    fallback = {
        "score": 0,
        "level": "revise",
        "issues": ["最终审稿结果解析失败。"],
        "revision_guidance": clean_text(review_text)[:500],
        "chinese_final": "",
        "english_final": "",
    }
    if not review_text:
        return fallback
    cleaned = review_text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    match = re.search(r"\{.*\}", cleaned, flags=re.DOTALL)
    if match:
        cleaned = match.group(0)
    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        return fallback
    issues = data.get("issues") or []
    if isinstance(issues, str):
        issues = [issues]
    try:
        score = int(data.get("score", 0))
    except (TypeError, ValueError):
        score = 0
    level = str(data.get("level", "revise")).strip().lower()
    if level not in {"excellent", "usable", "revise"}:
        level = "usable" if score >= 80 else "revise"
    return {
        "score": max(0, min(100, score)),
        "level": level,
        "issues": [clean_text(str(item)) for item in issues if clean_text(str(item))],
        "revision_guidance": clean_text(str(data.get("revision_guidance", ""))),
        "chinese_final": clean_summary_output(str(data.get("chinese_final", ""))),
        "english_final": clean_english_summary_output(str(data.get("english_final", ""))),
    }


def format_top_venue_final_report(review: dict) -> str:
    """把最终摘要审稿结果整理为报告文本。"""
    issues = review.get("issues") or []
    issue_text = "；".join(issues) if issues else "未发现明显结构或表达问题。"
    return "\n".join(
        [
            "Final Review：顶刊/顶会摘要自动审稿与中英文改写已完成。",
            f"Final Review：评分{review.get('score', 0)}/100，等级{review.get('level', 'revise')}。",
            f"Final Review：主要问题：{issue_text}",
            f"Final Review：修订建议：{review.get('revision_guidance', '') or '已按顶刊/顶会摘要结构完成改写。'}",
        ]
    )


def apply_top_venue_final_review(api_key: str, api_base: str, models: list[str], source_text: str, summary: str, summary_mode: str, cn_target: int | None = None, en_target: int | None = None) -> tuple[str, str, str]:
    """学术摘要最终审稿：自动给出中文和英文终稿，并返回可展示报告。"""
    review_text, used_model = call_llm_summary(
        api_key,
        api_base,
        models,
        top_venue_final_review_prompt(source_text, summary, summary_mode, cn_target, en_target),
        max_tokens=1200,
        timeout=60,
    )
    review = parse_top_venue_final_review(review_text)
    chinese_final = clean_summary_output(review.get("chinese_final", ""))
    english_final = clean_english_summary_output(review.get("english_final", ""))
    if not chinese_final:
        chinese_final = summary
    if summary_mode == "英文摘要":
        final_summary = english_final or clean_english_summary_output(summary)
    elif summary_mode in {"双语摘要", "中文摘要"}:
        final_summary = f"{chinese_final}\n\nEnglish Summary\n{english_final}" if english_final else chinese_final
    else:
        final_summary = chinese_final
    return final_summary.strip(), used_model, format_top_venue_final_report(review)


def pdca_revise_summary_prompt(source_text: str, summary: str, review: dict, academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """根据 PDCA 审查意见构建 Act 修订提示词。"""
    output_rule = summary_output_spec(summary_mode, academic, material_style, cn_target, en_target)
    issues = "；".join(review.get("issues") or ["无明确问题，但需按标准复核忠实性和表达。"])
    guidance = review.get("revision_guidance") or "删除无依据内容，补足原文支持的关键条件，保持客观克制。"
    material_label = prompt_material_label(academic, material_style)
    return f"""请执行PDCA循环中的Act环节：根据Check审查意见修订摘要。

Plan质量标准：
{build_pdca_plan(academic, summary_mode, material_style, cn_target, en_target)}

材料类型：
{material_label}

Do阶段初始摘要：
{summary}

Check阶段发现的问题：
{issues}

修订指引：
{guidance}

修订要求：
1. {output_rule}
2. 只保留原文支持的事实，不新增数据、指标、模型名、结论或建议。
3. 删除或弱化“证明、显著、最佳、最优、企业级部署、全面解决”等缺少充分依据的表达。
4. 原文没有定量结果时，只写定性发现；原文没有限制条件时，不主动编造局限。
5. 区分技术贡献、实验结果和作者主张，避免把推断写成事实。
6. 如果是PPT或图片，按材料主题、内容结构、关键要点、阅读提示来组织，不套用论文或业务五段式结论。
7. 如果是双语摘要，中文摘要和 English Summary 必须分开输出，不能同段混排。
8. 如果是学术材料，应修成顶刊/顶会摘要风格：以背景缺口或现实问题开篇，承接研究任务和方法设计，再给出原文支持的证据、贡献和含义；不要写成读书笔记或文档介绍。
9. 删除所有非摘要正文内容，包括JSON字段、评分、审稿意见、提示词残留、系统说明、对用户的操作建议、“以下是”“根据你的要求”等开场白。
10. 逐句对照原文证据：没有直接依据的句子必须删除；重要但证据不足的信息只能写成“原文未说明”或更保守的限定表述。
11. 数字、样本量、模型名称、指标、比较结论、因果关系和贡献声明必须来自原文；原文没有出现时不能补写。
12. 只输出最终摘要，不输出证据说明、修改过程、评分或质量检查报告。

原文：
{source_text[:9000]}

请直接输出修订后的最终摘要。
"""


def build_pdca_report(plan: str, cycle_records: list[dict]) -> str:
    """把多轮 PDCA 的计划、审查和修订动作整理为可展示的检查报告。"""
    lines = [f"Plan：{plan}", "Do：已生成初始摘要。"]
    for record in cycle_records:
        review = record["review"]
        issues = review.get("issues") or []
        issue_text = "；".join(issues) if issues else "未发现明显幻觉、夸大或关键条件遗漏。"
        action = (
            "已根据审查意见完成Act修订，并进入下一轮Check。"
            if record.get("revised")
            else "Check通过，保留当前摘要。"
        )
        if record.get("stopped_at_limit"):
            action = "已达到最大修订轮次，保留最后一版摘要，建议人工复核。"
        lines.append(
            f"Cycle {record['round']}：Check评分{review.get('score', 0)}/100，"
            f"结论{review.get('verdict', 'revise')}。{issue_text}"
        )
        lines.append(f"Cycle {record['round']}：Act：{action}")
    return "\n".join(lines)


def apply_pdca_summary_cycle(api_key: str, api_base: str, models: list[str], source_text: str, summary: str, academic: bool, max_tokens: int = 750, summary_mode: str = "中文摘要", max_cycles: int = 3, pass_score: int = 88, material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> tuple[str, str, str]:
    """执行多轮摘要 PDCA 质量循环：审查、修订、复查，直到通过或达到轮次上限。"""
    plan = build_pdca_plan(academic, summary_mode, material_style, cn_target, en_target)
    used_model = models[0] if models else ""
    cycle_records = []
    current_summary = clean_summary_output(summary)

    for round_index in range(1, max_cycles + 1):
        review_text, review_model = call_llm_summary(
            api_key,
            api_base,
            models,
            pdca_check_summary_prompt(source_text, current_summary, academic, summary_mode, material_style, cn_target, en_target),
            max_tokens=550,
            timeout=45,
        )
        used_model = review_model
        review = parse_pdca_review(review_text, pass_score=pass_score)
        top_venue_style_gap = academic and material_style == "document" and needs_top_venue_polish(current_summary, source_text)
        if top_venue_style_gap:
            review["verdict"] = "revise"
            review["score"] = min(review.get("score", 0), pass_score - 1)
            review.setdefault("issues", [])
            review["issues"].append("摘要尚未达到顶刊/顶会式研究摘要结构，需要强化背景缺口、方法设计、证据结果和边际贡献之间的承接。")
            review["revision_guidance"] = (
                "请按顶刊/顶会摘要结构重写：先交代研究缺口或现实需求，再说明任务重构和方法机制，"
                "随后给出原文支持的数据/实验结果，最后收束到边际贡献和应用含义。"
            )
        should_revise = (
            review.get("verdict") == "revise"
            or bool(review.get("issues"))
            or review.get("score", 0) < pass_score
            or is_low_quality_summary(current_summary, academic=academic, summary_mode=summary_mode)
            or needs_style_polish(current_summary)
            or top_venue_style_gap
        )
        record = {"round": round_index, "review": review, "revised": False, "stopped_at_limit": False}
        cycle_records.append(record)
        if not should_revise:
            break
        if round_index >= max_cycles:
            record["stopped_at_limit"] = True
            break
        current_summary, used_model = call_llm_summary(
            api_key,
            api_base,
            models,
            pdca_revise_summary_prompt(source_text, current_summary, review, academic, summary_mode, material_style, cn_target, en_target),
            max_tokens=max_tokens,
            timeout=45,
        )
        current_summary = clean_summary_output(current_summary)
        record["revised"] = True

    report = build_pdca_report(plan, cycle_records)
    return current_summary, used_model, report


def critique_summary_prompt(source_text: str, summary: str, academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """构建风格润色提示词，用于把机械摘要改成更自然的中文表达。"""
    standard = summary_output_spec(summary_mode, academic, material_style, cn_target, en_target)
    return f"""请根据原文重写下面的摘要，让它更像人工写作，而不是关键词压缩或概念堆叠。

标准：{standard}

改写要求：
1. 只输出改写后的中文摘要，不要解释。
2. 保留原文支持的事实，不新增原文没有的信息。
3. 拆开过长句，减少名词堆叠，让句子有自然的主谓宾。
4. 避免连续使用“围绕、系统梳理、进一步分析、归纳了、研究表明、相关内容为”等模板表达。
5. 如果原摘要像清单，请改成有层次的自然表达。
6. 删除箭头、加号、斜杠和括号式解释，把它们改写成自然中文句子。
7. 删除“原文说明、原文指出、该文档提到、本文旨在、这段内容主要讨论”等转述腔，改成直接陈述研究对象、事项、方法、发现或结论。

原文：
{source_text[:9000]}

当前摘要：
{summary}
"""


def build_chunk_summary_prompt(chunk: str, index: int, total: int, academic: bool, summary_mode: str = "中文摘要", material_style: str = "document", cn_target: int | None = None, en_target: int | None = None) -> str:
    """为长文分块摘要构建单个片段的事实抽取提示词。"""
    language_rule = "只输出中文。" if summary_mode == "中文摘要" else "Use English for extracted facts." if summary_mode == "英文摘要" else "同时保留中文和英文可用事实。"
    if material_style in {"presentation", "visual"}:
        material_name = "PPT片段" if material_style == "presentation" else "图片OCR片段"
        return f"""请从以下{material_name}中抽取用于最终摘要的事实。
要求：
1. {language_rule}
2. 用 4-6 条短句列出材料主题、页面结构、关键机制/流程/指标/变量和可读设定。
3. 不要写最终摘要，不要强行补写实验结论、风险或行动建议。
4. 如果本片段只是标题或图中文字，请保持克制，只记录可确认信息。
5. 不要编造数字、样本量、年份或效果判断。

片段 {index}/{total}：
{chunk}
"""
    if academic:
        return f"""请从以下论文片段中抽取用于最终中文学术摘要的事实。
要求：
1. {language_rule}
2. 用 4-6 条短句列出事实，不要写最终摘要。
3. 按顶刊/顶会摘要需要，优先保留背景缺口、研究问题、方法/模型、数据或实验、主要结果、贡献边界和局限。
4. 保留必要英文术语。
5. 不要编造数字、样本量、年份或结论。
6. 如果本片段信息重复，只保留新增事实。

论文片段 {index}/{total}：
{chunk}
"""

    return f"""请从以下文档片段中抽取用于最终摘要的事实。
要求：
1. {language_rule}
2. 用 4-6 条短句列出事实，不要写最终摘要。
3. 保留实体、日期、金额、结论、风险、要求和行动项。
4. 不要编造原文没有的信息。
5. 如果本片段信息重复，只保留新增事实。

文档片段 {index}/{total}：
{chunk}
"""


def build_llm_source(text: str, academic: bool, limit: int = 9000) -> str:
    """从长文中截取最适合送入大模型的代表性内容。"""
    text = text.strip()
    if len(text) <= limit:
        return text

    if not academic:
        head = text[: limit // 2]
        tail = text[-(limit // 2):]
        return clean_text(f"{head}\n\n……\n\n{tail}")[:limit]

    section_patterns = [
        r"(?:^|\n)\s*(?:摘要|abstract)\s*[:：]?\s*(.*?)(?=\n\s*(?:关键词|keywords|引言|introduction|1[.\s]|一、)|$)",
        r"(?:^|\n)\s*(?:引言|introduction)\s*[:：]?\s*(.*?)(?=\n\s*(?:方法|method|methodology|实验|experiment|结果|results|2[.\s]|二、)|$)",
        r"(?:^|\n)\s*(?:结论|conclusion|conclusions)\s*[:：]?\s*(.*?)(?=\n\s*(?:参考文献|references|致谢|acknowledg|$))",
    ]
    parts = []
    for pattern in section_patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE | re.DOTALL)
        if match:
            parts.append(match.group(0)[:2500])

    parts.append(text[:3500])
    parts.append(text[-2500:])
    source = clean_text("\n\n".join(parts))
    return source[:limit]


def summarize_chunks_with_llm(api_key: str, api_base: str, models: list[str], text: str, academic: bool, summary_mode: str = "中文摘要", pdca_cycles: int = 3, pdca_pass_score: int = 88, material_style: str = "document", cn_target: int | None = None, en_target: int | None = None, quality_mode: str = "快速") -> tuple[str, str, str]:
    """长文摘要流程：先分块抽取事实，再合并生成最终摘要。"""
    chunks = split_text_chunks(text, chunk_size=3500, overlap=250)
    partial_summaries = []
    used_model = models[0]
    for index, chunk in enumerate(chunks, start=1):
        partial, partial_model = call_llm_summary(
            api_key,
            api_base,
            models,
            build_chunk_summary_prompt(chunk, index, len(chunks), academic, summary_mode, material_style, cn_target, en_target),
            max_tokens=450,
            timeout=30,
        )
        used_model = partial_model
        if partial and len(partial.strip()) >= 20 and not looks_mojibake(partial):
            partial_summaries.append(partial)

    if not partial_summaries:
        raise RuntimeError("empty chunk summaries")

    combined = "\n\n".join(partial_summaries)
    summary, used_model = call_llm_summary(api_key, api_base, models, build_summary_prompt(combined, academic, summary_mode, material_style, cn_target, en_target), max_tokens=750)
    pdca_report = ""
    if summary:
        summary, used_model, pdca_report = apply_pdca_summary_cycle(
            api_key,
            api_base,
            models,
            combined,
            summary,
            academic,
            max_tokens=750,
            summary_mode=summary_mode,
            max_cycles=pdca_cycles,
            pass_score=pdca_pass_score,
            material_style=material_style,
            cn_target=cn_target,
            en_target=en_target,
        )
        if quality_mode == "精修" and academic and material_style == "document":
            try:
                summary, review_model, final_report = apply_top_venue_final_review(
                    api_key, api_base, models, combined, summary, summary_mode, cn_target, en_target
                )
                used_model = review_model
                pdca_report = "\n".join([item for item in [pdca_report, final_report] if item])
            except Exception as review_exc:
                pdca_report = "\n".join(
                    [
                        pdca_report,
                        f"Final Review：顶刊/顶会摘要最终审稿暂未完成，保留PDCA摘要。原因：{str(review_exc)[:120]}",
                    ]
                )
    if not summary or is_low_quality_summary(summary, academic=academic, summary_mode=summary_mode):
        raise RuntimeError("low quality final summary")
    return summary, used_model, pdca_report


def summarize_text_with_llm(text: str, summary_mode: str = "中文摘要", pdca_cycles: int = 3, pdca_pass_score: int = 88, material_style: str = "document", cn_target: int | None = None, en_target: int | None = None, quality_mode: str = "快速") -> tuple[str, str, str]:
    """摘要总入口：优先大模型和 PDCA，失败后回退本地摘要。"""
    api_key, api_base, models = get_llm_config()
    academic = is_academic_document(text)
    if not api_key:
        if summary_mode != "中文摘要":
            local_model_result, local_model_error = local_model_summary(text, summary_mode, academic)
            if local_model_result:
                return local_model_result, f"本地{summary_mode}模型", "未调用在线质量复核；本地模型输出建议人工复核。"
            local = summarize_text(text, summary_mode=summary_mode, material_style=material_style)
            return local, "本地摘要", f"未配置在线摘要服务，且本地英文/双语模型不可用：{local_model_error}"
        return summarize_text(text, summary_mode=summary_mode, material_style=material_style), "本地摘要", "未配置大模型API，使用本地摘要兜底，未执行PDCA对抗审查。"

    if is_unreadable_text(text):
        return (
            "提取文本疑似存在大量乱码或不可读字符，系统已停止生成正式摘要。请重新导出文件、检查编码或上传更清晰的原始文档。",
            "文本质量检查",
            "文本质量检查未通过，未执行摘要PDCA循环。",
        )

    try:
        source = build_llm_source(text, academic)
        summary, used_model = call_llm_summary(api_key, api_base, models, build_summary_prompt(source, academic, summary_mode, material_style, cn_target, en_target), max_tokens=700)
        pdca_report = ""
        if summary:
            initial_summary = clean_summary_output(summary)
            try:
                summary, used_model, pdca_report = apply_pdca_summary_cycle(
                    api_key,
                    api_base,
                    models,
                    source,
                    initial_summary,
                    academic,
                    max_tokens=700,
                    summary_mode=summary_mode,
                    max_cycles=pdca_cycles,
                    pass_score=pdca_pass_score,
                    material_style=material_style,
                    cn_target=cn_target,
                    en_target=en_target,
                )
                if quality_mode == "精修" and academic and material_style == "document":
                    try:
                        summary, review_model, final_report = apply_top_venue_final_review(
                            api_key, api_base, models, source, summary, summary_mode, cn_target, en_target
                        )
                        used_model = review_model
                        pdca_report = "\n".join([item for item in [pdca_report, final_report] if item])
                    except Exception as review_exc:
                        pdca_report = "\n".join(
                            [
                                pdca_report,
                                f"Final Review：顶刊/顶会摘要最终审稿暂未完成，保留PDCA摘要。原因：{str(review_exc)[:120]}",
                            ]
                        )
            except Exception as pdca_exc:
                summary = initial_summary
                pdca_report = f"质量复核暂未完成，已保留初始摘要。原因：{str(pdca_exc)[:120]}"
        if summary and not is_low_quality_summary(summary, academic=academic, summary_mode=summary_mode):
            if len(text) > len(source):
                source_name = "智能学术摘要（长文代表片段）" if academic else "智能摘要（长文代表片段）"
            else:
                source_name = "智能学术摘要" if academic else "智能摘要"
            return summary, f"{source_name} + 质量复核优化", pdca_report

        if len(text) > 6000:
            summary, used_model, pdca_report = summarize_chunks_with_llm(
                api_key, api_base, models, text, academic, summary_mode,
                pdca_cycles=pdca_cycles, pdca_pass_score=pdca_pass_score, material_style=material_style,
                cn_target=cn_target, en_target=en_target, quality_mode=quality_mode
            )
            source_name = "智能学术分块摘要" if academic else "智能分块摘要"
            return summary, f"{source_name} + 质量复核优化", pdca_report
    except Exception as exc:
        if summary_mode != "中文摘要":
            local_model_result, local_model_error = local_model_summary(text, summary_mode, academic)
            if local_model_result:
                return local_model_result, f"本地{summary_mode}模型（在线摘要不可用：{str(exc)[:60]}）", "未完成在线PDCA质量复核；本地模型输出建议人工复核。"
        local_summary = summarize_text(text, summary_mode=summary_mode, material_style=material_style)
        return local_summary, f"本地摘要（智能摘要服务暂不可用：{str(exc)[:80]}）", "智能摘要服务暂不可用，已回退本地摘要，未完成PDCA对抗审查。"

    return summarize_text(text, summary_mode=summary_mode, material_style=material_style), "本地摘要", "未生成可用大模型摘要，使用本地摘要兜底，未执行PDCA对抗审查。"

def classify_document(text: str) -> tuple[str, dict[str, int]]:
    """基于关键词规则给文档分类，并返回各类别得分。"""
    scores = {
        category: sum(text.count(keyword) for keyword in keywords)
        for category, keywords in STANDARD_CATEGORY_RULES.items()
    }
    best_category = max(scores, key=scores.get)
    if scores[best_category] == 0:
        return "通用文档", scores
    return best_category, scores


def build_category_explanation(category: str, category_scores: dict[str, int], text: str = "") -> str:
    """生成自动分类依据说明，便于页面和导出报告展示验收点。"""
    sorted_scores = sorted(category_scores.items(), key=lambda item: item[1], reverse=True)
    top_scores = "；".join(f"{name}={score}" for name, score in sorted_scores[:5])
    matched_keywords = []
    for keyword in STANDARD_CATEGORY_RULES.get(category, []):
        if keyword and keyword in text:
            matched_keywords.append(keyword)
        if len(matched_keywords) >= 8:
            break
    if matched_keywords:
        evidence = "命中关键词：" + "、".join(matched_keywords)
    elif category == "通用文档":
        evidence = "未命中明显领域关键词，归入通用文档。"
    else:
        evidence = "基于类别关键词累计得分最高判定。"
    return f"系统按预设类别关键词对全文进行累计打分，得分最高类别为“{category}”。{evidence} Top得分：{top_scores}。"


def resolve_user_category(auto_category: str) -> str:
    """根据侧边栏设置决定最终保存分类：自动、标准分类或用户自定义分类。"""
    mode = st.session_state.get("category_mode", "自动分类")
    if mode == "选择标准分类":
        manual_category = clean_text(str(st.session_state.get("manual_category", "")))
        return manual_category or auto_category
    if mode == "自定义分类":
        custom_category = clean_text(str(st.session_state.get("custom_category", "")))
        return custom_category or auto_category
    return auto_category


def build_notice_figure(message: str):
    """生成稳定的占位图，避免词云异常时页面或导出中断。"""
    figure, axis = plt.subplots(figsize=(11, 5))
    axis.text(0.5, 0.5, message, ha="center", va="center", fontsize=16)
    axis.axis("off")
    figure.tight_layout(pad=0.2)
    return figure


def resolve_cjk_font_path() -> str | None:
    """查找本机或云端常见中文字体，供词云和图片导出使用。"""
    font_candidates = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    return next((path for path in font_candidates if Path(path).exists()), None)


def build_wordcloud(freq_df: pd.DataFrame, top_n: int = 40):
    """根据词频表生成中文词云图。"""
    word_col = "词语" if "词语" in freq_df.columns else "word"
    count_col = "频次" if "频次" in freq_df.columns else "count"
    filtered = freq_df.copy()
    if filtered.empty or word_col not in filtered.columns or count_col not in filtered.columns:
        return build_notice_figure("暂无足够词频生成词云")
    filtered[word_col] = filtered[word_col].astype(str).str.strip()
    filtered = filtered[filtered[word_col].astype(bool)]
    filtered = filtered[~filtered[word_col].isin(STOPWORDS)]
    filtered[count_col] = pd.to_numeric(filtered[count_col], errors="coerce").fillna(0)
    filtered = filtered[filtered[count_col] > 0]
    preferred = filtered[filtered[count_col] >= 2]
    if not preferred.empty:
        filtered = preferred
    filtered = filtered.sort_values(count_col, ascending=False).head(top_n)
    frequencies = dict(zip(filtered[word_col], filtered[count_col]))
    if not frequencies:
        return build_notice_figure("暂无足够词频生成词云")
    try:
        wordcloud = WordCloud(
            font_path=resolve_cjk_font_path(),
            width=1100,
            height=520,
            background_color="#fbfcfe",
            colormap="cividis",
            max_words=min(top_n, 45),
            prefer_horizontal=1.0,
            collocations=False,
            margin=10,
            min_font_size=10,
            max_font_size=92,
            relative_scaling=0.45,
            contour_width=1,
            contour_color="#e5e7eb",
            random_state=42,
        ).generate_from_frequencies(frequencies)
    except Exception as exc:
        return build_notice_figure(f"词云生成失败，已保留词频柱状图\n{str(exc)[:90]}")

    figure, axis = plt.subplots(figsize=(11, 5), facecolor="#fbfcfe")
    axis.imshow(wordcloud, interpolation="bilinear")
    axis.axis("off")
    figure.tight_layout(pad=0.4)
    return figure


def build_word_frequency_chart(freq_df: pd.DataFrame, top_n: int = 30):
    """构建词频柱状图，供页面展示和 HTML 下载复用。"""
    word_col = "词语" if "词语" in freq_df.columns else "word"
    count_col = "频次" if "频次" in freq_df.columns else "count"
    chart_df = freq_df.sort_values(count_col, ascending=False).head(top_n)
    fig = px.bar(chart_df, x=word_col, y=count_col, text=count_col, title=f"高频词 Top {top_n}")
    fig.update_layout(xaxis_title="词语", yaxis_title="频次")
    return fig


def figure_to_png_bytes(figure) -> bytes:
    """把 Matplotlib 图导出为 PNG 字节。"""
    buffer = io.BytesIO()
    figure.savefig(buffer, format="png", dpi=180, bbox_inches="tight")
    plt.close(figure)
    return buffer.getvalue()


def parse_uploaded_file(uploaded_file, ocr_mode: str = "standard", pdf_extract_mode: str = "auto") -> tuple[str, bool]:
    """按文件后缀选择 PDF、Word、PPT、表格、文本或图片解析方式。"""
    suffix = uploaded_file.name.lower().rsplit(".", 1)[-1]
    file_bytes = uploaded_file.getvalue()
    validate_file_kind(file_bytes, uploaded_file.name)
    file_hash = hashlib.sha256(file_bytes).hexdigest()
    parser_options = json.dumps(
        {"ocr_mode": ocr_mode, "pdf_extract_mode": pdf_extract_mode, "suffix": suffix},
        ensure_ascii=False,
        sort_keys=True,
    )
    cache_key = build_parse_cache_key(file_hash, uploaded_file.name, ocr_mode, pdf_extract_mode)
    cached_text = get_parse_cache(cache_key)
    if cached_text:
        return cached_text, True

    if suffix == "pdf":
        extracted = post_process_extracted_text(
            extract_pdf(file_bytes, ocr_mode=ocr_mode, pdf_extract_mode=pdf_extract_mode, filename=uploaded_file.name),
            "pdf",
        )
    elif suffix == "docx":
        extracted = post_process_extracted_text(extract_docx(io.BytesIO(file_bytes)), "document")
    elif suffix == "doc":
        extracted = post_process_extracted_text(extract_doc(file_bytes, uploaded_file.name), "document")
    elif suffix == "pptx":
        extracted = post_process_extracted_text(extract_pptx(io.BytesIO(file_bytes)), "ppt")
    elif suffix == "ppt":
        extracted = post_process_extracted_text(extract_ppt(file_bytes, uploaded_file.name), "ppt")
    elif suffix in {"csv", "xlsx", "xls"}:
        extracted = post_process_extracted_text(extract_tabular(io.BytesIO(file_bytes), suffix), "table")
    elif suffix in {"txt", "md", "json", "jsonl"}:
        extracted = post_process_extracted_text(extract_plain_text(file_bytes), "document")
    elif suffix in {"png", "jpg", "jpeg", "bmp", "tif", "tiff"}:
        extracted = post_process_extracted_text(extract_image(io.BytesIO(file_bytes), ocr_mode=ocr_mode), "image")
    else:
        raise ValueError("暂不支持该文件格式")

    save_parse_cache(cache_key, file_hash, uploaded_file.name, parser_options, extracted)
    return extracted, False


def analyze_text(text: str, filename: str = "", summary_mode: str = "中文摘要", pdca_cycles: int = 3, pdca_pass_score: int = 88, cn_target: int | None = None, en_target: int | None = None, quality_mode: str = "快速") -> tuple[str, dict[str, int], str, str, str, pd.DataFrame, dict]:
    """对提取文本执行分类、摘要和词频统计。"""
    route = route_file(filename, text).to_dict()
    material_style = infer_material_style(filename, route.get("file_type", ""))
    category, category_scores = classify_document(text)
    summary, summary_source, pdca_report = summarize_text_with_llm(
        text,
        summary_mode=summary_mode,
        pdca_cycles=pdca_cycles,
        pdca_pass_score=pdca_pass_score,
        material_style=material_style,
        cn_target=cn_target,
        en_target=en_target,
        quality_mode=quality_mode,
    )
    chinese_summary, english_summary = split_bilingual_summary(summary)
    if english_summary:
        summary = f"{clean_summary_output(chinese_summary)}\n\nEnglish Summary\n{clean_english_summary_output(english_summary)}"
    else:
        summary = clean_summary_output(summary)
    freq = Counter(tokenize(text)).most_common(30)
    freq_df = pd.DataFrame(freq, columns=["word", "count"])
    return category, category_scores, summary, summary_source, pdca_report, freq_df, route


def render_route(route: dict | None) -> None:
    """展示文件类型识别、任务路由和推荐输出。"""
    if not route:
        return
    st.subheader("文件路由与任务识别")
    c1, c2, c3 = st.columns([1, 1, 1])
    c1.metric("文件类型", route.get("file_type_label", "未知"))
    c2.metric("任务类别", route.get("task_label", "通用文档"))
    c3.metric("置信度", f"{float(route.get('confidence', 0)):.0%}")
    st.caption(route.get("reason", ""))
    outputs = route.get("recommended_outputs") or []
    if outputs:
        st.write("推荐输出：" + "、".join(outputs))


def split_bilingual_summary(summary: str) -> tuple[str, str]:
    """把中英双语摘要拆成中文和英文两部分。"""
    text = clean_text(summary)
    if not text:
        return "", ""
    marker_match = re.search(r"(?:English Summary|English summary|英文摘要)\s*[:：]?", text)
    if marker_match and marker_match.start() > 20:
        chinese_part = text[: marker_match.start()].strip(" \n。")
        english_part = text[marker_match.end():].strip()
        if len(re.findall(r"[\u4e00-\u9fff]", chinese_part)) >= 20 and len(re.findall(r"[A-Za-z]{3,}", english_part)) >= 20:
            return chinese_part, english_part

    english_leads = (
        "In this study",
        "This study",
        "This paper",
        "The paper",
        "We examine",
        "We study",
        "Financial distress",
        "Corporate governance",
        "Bankruptcy prediction",
        "Credit risk",
    )
    for index, char in enumerate(text):
        if char not in "。！？；\n":
            continue
        start = index + 1
        while start < len(text) and text[start].isspace():
            start += 1
        candidate = text[start : start + 240]
        starts_like_english = candidate.startswith(english_leads) or bool(
            re.match(r"[A-Z][A-Za-z]+(?:[ ,'-]+[A-Za-z]+){4,}", candidate)
        )
        if starts_like_english:
            chinese_part = text[:start].strip(" \n。")
            english_part = text[start:].strip()
            if len(re.findall(r"[\u4e00-\u9fff]", chinese_part)) >= 20 and len(re.findall(r"[A-Za-z]{3,}", english_part)) >= 20:
                return chinese_part, english_part

    sentences = re.split(r"(?<=[。！？.!?])\s*", text)
    chinese_parts = []
    english_parts = []
    for sentence in sentences:
        if not sentence.strip():
            continue
        chinese_count = len(re.findall(r"[\u4e00-\u9fff]", sentence))
        english_count = len(re.findall(r"[A-Za-z]", sentence))
        if english_count > chinese_count * 1.5 and english_count > 20:
            english_parts.append(sentence.strip())
        else:
            chinese_parts.append(sentence.strip())
    if chinese_parts and english_parts:
        return " ".join(chinese_parts).strip(), " ".join(english_parts).strip()
    return text, ""


def count_chinese_and_english(text: str) -> tuple[int, int]:
    """统计中文字符数和英文单词数，用于页面展示更直观的篇幅信息。"""
    chinese_count = len(re.findall(r"[\u4e00-\u9fff]", text))
    english_word_count = len(re.findall(r"[A-Za-z]+(?:'[A-Za-z]+)?", text))
    return chinese_count, english_word_count


def render_summary(summary: str, summary_source: str) -> None:
    """展示摘要；双语内容自动拆成两个框。"""
    st.subheader("文本摘要")
    st.caption(f"摘要来源：{summary_source}")
    chinese_summary, english_summary = split_bilingual_summary(summary)
    if english_summary:
        col_cn, col_en = st.columns(2)
        with col_cn:
            st.markdown("**中文摘要**")
            st.info(chinese_summary)
        with col_en:
            st.markdown("**English Summary**")
            st.info(english_summary)
        return
    st.info(summary)


def format_extracted_text_for_reading(text: str, max_chars: int = 12000) -> str:
    """把提取文本整理成适合页面阅读的段落预览。"""
    cleaned = clean_text(text)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n", cleaned) if item.strip()]
    if len(paragraphs) <= 1:
        sentences = split_sentences(cleaned)
        paragraphs = ["".join(sentences[index:index + 4]) for index in range(0, len(sentences), 4)]
    preview = "\n\n".join(paragraphs)
    if len(preview) > max_chars:
        preview = preview[:max_chars].rstrip() + "\n\n……后续内容已折叠，可在下方查看原始全文。"
    return preview or "暂无可展示文本。"


def extracted_text_quality(text: str) -> tuple[str, str, list[str]]:
    """评估全文提取质量，给普通用户可执行的处理建议。"""
    cleaned = clean_text(text)
    compact = re.sub(r"\s+", "", cleaned)
    if not compact:
        return "低", "未提取到有效文本", ["请重新上传清晰 PDF/Word，扫描件建议选择“扫描件识别”。"]

    lines = [line.strip() for line in cleaned.splitlines() if line.strip()]
    short_ratio = sum(1 for line in lines if len(line) <= 8) / max(len(lines), 1)
    readable_ratio = len(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]", compact)) / max(len(compact), 1)
    replacement_char = chr(0xFFFD)
    noise_ratio = (
        compact.count(replacement_char)
        + sum(compact.count(char) for char in "□■●◆◇")
        + mojibake_marker_count(compact)
    ) / max(len(compact), 1)

    issues = []
    score = 100
    formula_line_count = sum(1 for line in lines if looks_like_formula_line(line))
    table_line_count = sum(1 for line in lines if line.startswith("|") and line.endswith("|"))
    if is_unreadable_text(cleaned) or noise_ratio > 0.04:
        score -= 45
        issues.append("疑似存在乱码或字体编码问题，建议重新导出 PDF，或上传原始 Word。")
    if formula_line_count:
        issues.append(f"检测到 {formula_line_count} 行疑似公式或数学表达式；系统会尽量保留公式结构，复杂公式建议使用“多模态精读”或高精度 OCR 后人工复核。")
    if table_line_count:
        issues.append(f"检测到 {table_line_count} 行表格结构；系统已尽量按 Markdown 表格保留列和行，复杂合并单元格建议人工复核。")
    if "图片结构化识别" in cleaned:
        issues.append("检测到图片结构化识别结果；图表、流程图和图片文字已按可读内容转写，关键数值建议对照原图复核。")
    if short_ratio > 0.45 and len(lines) >= 8:
        score -= 25
        issues.append("短行比例偏高，可能是扫描件、分栏 PDF 或逐字断行，建议尝试“多模态精读”或“扫描件识别”。")
    if readable_ratio < 0.45:
        score -= 20
        issues.append("可读字符比例偏低，可能包含大量图片、公式或表格截图。")
    if len(compact) < 500:
        score -= 15
        issues.append("提取文本较短，摘要可能不完整。")

    if score >= 80:
        return "较好", "文本结构基本可读", issues or ["当前提取结果可直接用于摘要，正式提交前仍建议人工复核。"]
    if score >= 55:
        return "一般", "可用但需要复核", issues
    return "偏低", "提取质量可能影响摘要", issues


def render_extracted_text(text: str) -> None:
    """展示排版后的全文预览，并保留原始文本方便复制。"""
    st.subheader("提取的全文")
    cn_count, en_count = count_chinese_and_english(text)
    quality_label, quality_note, quality_issues = extracted_text_quality(text)
    c1, c2, c3 = st.columns(3)
    c1.metric("中文字符", cn_count)
    c2.metric("英文单词", en_count)
    c3.metric("总字符", len(text))
    if quality_label == "较好":
        st.success(f"提取质量：{quality_label}｜{quality_note}")
    elif quality_label == "一般":
        st.warning(f"提取质量：{quality_label}｜{quality_note}")
    else:
        st.error(f"提取质量：{quality_label}｜{quality_note}")
    with st.expander("提取质量建议", expanded=quality_label != "较好"):
        for issue in quality_issues:
            st.write(f"- {issue}")
    st.markdown(format_extracted_text_for_reading(text).replace("\n", "\n\n"))
    with st.expander("查看原始文本", expanded=False):
        st.text_area("原始全文", text, height=360)


def render_analysis(category: str, category_scores: dict[str, int], summary: str, summary_source: str, pdca_report: str, freq_df: pd.DataFrame, extracted_text: str, route: dict | None = None) -> None:
    """在 Streamlit 页面中渲染分类、摘要、词频图和全文。"""
    render_route(route)
    left, right = st.columns([1, 1])

    with left:
        st.subheader("自动分类")
        st.success(category)
        st.caption("分类方式：基于预设类别关键词对全文累计打分，得分最高者作为自动类别。")
        with st.expander("查看自动分类依据", expanded=False):
            st.write(build_category_explanation(category, category_scores, extracted_text))

    with right:
        st.subheader("文档信息")
        st.metric("文本长度", f"{len(extracted_text)} 字符")
        st.metric("关键词数量", len(freq_df))

    render_summary(summary, summary_source)
    if pdca_report:
        with st.expander("PDCA对抗式摘要质量检查", expanded=False):
            st.text(pdca_report)

    st.subheader("词频统计")
    if freq_df.empty:
        st.write("暂无可统计词语。")
    else:
        tab_bar, tab_cloud, tab_class = st.tabs(["柱状图", "词云", "分类得分"])
        word_col = "词语" if "词语" in freq_df.columns else "word"
        count_col = "频次" if "频次" in freq_df.columns else "count"

        with tab_bar:
            fig = build_word_frequency_chart(freq_df)
            st.plotly_chart(fig, width="stretch")

        with tab_cloud:
            try:
                st.pyplot(build_wordcloud(freq_df), clear_figure=True)
            except Exception as exc:
                st.write(f"词云生成失败：{exc}。柱状图仍可正常使用。")

        with tab_class:
            score_df = pd.DataFrame(
                sorted(category_scores.items(), key=lambda item: item[1], reverse=True),
                columns=["类别", "得分"],
            )
            st.dataframe(score_df, width="stretch")

    render_extracted_text(extracted_text)

def estimate_processing_time(filename: str, file_size: int, ocr_mode: str, summary_mode: str) -> str:
    """按文件类型、大小和模式给用户一个保守的等待时间区间。"""
    suffix = Path(filename).suffix.lower()
    size_mb = file_size / (1024 * 1024)
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}:
        base_low, base_high = (8, 25) if ocr_mode == "fast" else (15, 45) if ocr_mode == "standard" else (30, 90)
    elif suffix == ".pdf":
        base_low, base_high = (10, 35) if size_mb < 5 else (25, 90) if size_mb < 20 else (60, 180)
    elif suffix in {".xlsx", ".xls", ".csv"}:
        base_low, base_high = (5, 20) if size_mb < 10 else (20, 60)
    else:
        base_low, base_high = (5, 20) if size_mb < 5 else (15, 50)
    if summary_mode in {"英文摘要", "双语摘要"}:
        base_low += 5
        base_high += 25
    if get_secret("LLM_API_KEY"):
        base_low += 10
        base_high += 50
    if base_high < 60:
        return f"预计约 {base_low}-{base_high} 秒"
    return f"预计约 {max(1, round(base_low / 60))}-{max(1, round(base_high / 60))} 分钟"


def update_progress(progress_bar, status_box, percent: int, message: str, eta_hint: str = "") -> None:
    """更新页面进度条，并给出保守耗时提示，降低长任务等待焦虑。"""
    progress_bar.progress(min(max(percent, 0), 100))
    if percent >= 100:
        status_box.success(f"{message} ｜ 进度 100%")
        return

    eta_text = f" ｜ {eta_hint}" if eta_hint and percent < 20 else ""
    status_box.info(f"{message} ｜ 进度 {percent}%{eta_text} ｜ 正在处理，请稍候")


def format_duration(seconds: float | int | None) -> str:
    """把处理秒数格式化为页面上容易理解的用时。"""
    try:
        seconds = float(seconds or 0)
    except (TypeError, ValueError):
        seconds = 0.0
    if seconds <= 0:
        return "未记录"
    if seconds < 60:
        return f"{seconds:.1f} 秒"
    minutes = int(seconds // 60)
    remain = int(round(seconds % 60))
    return f"{minutes} 分 {remain} 秒"


def safe_download_stem(filename: str) -> str:
    """生成下载文件名主干。"""
    stem = Path(filename).stem or "document"
    return re.sub(r"[^A-Za-z0-9_.\-\u4e00-\u9fff]", "_", stem)


def build_markdown_export(result: dict) -> str:
    """导出结构化 Markdown 分析报告。"""
    summary = str(result["summary"])
    chinese_summary, english_summary = split_bilingual_summary(summary)
    cn_chars, en_words = count_chinese_and_english(summary)
    text_cn_chars, text_en_words = count_chinese_and_english(str(result["extracted_text"]))
    freq_df = result.get("freq_df")
    freq_text = "暂无词频数据"
    if isinstance(freq_df, pd.DataFrame) and not freq_df.empty:
        word_col = "词语" if "词语" in freq_df.columns else "word"
        count_col = "频次" if "频次" in freq_df.columns else "count"
        freq_text = "、".join(f"{row[word_col]}({row[count_col]})" for _, row in freq_df.head(20).iterrows())
    summary_block = [f"## 中文摘要\n\n{chinese_summary}"]
    if english_summary:
        summary_block.append(f"## English Summary\n\n{english_summary}")
    return "\n\n".join(
        [
            "# 文档智能分析报告",
            f"**文件名：** {result['filename']}",
            f"**文档分类：** {result['category']}",
            f"**自动分类依据：** {build_category_explanation(str(result['category']), result.get('category_scores', {}), str(result['extracted_text']))}",
            f"**摘要统计：** 中文 {cn_chars} 字 / 英文 {en_words} 词",
            f"**原文统计：** 中文 {text_cn_chars} 字 / 英文 {text_en_words} 词",
            *summary_block,
            f"## 高频词\n\n{freq_text}",
            f"## 提取文本\n\n{format_extracted_text_for_reading(str(result['extracted_text']), max_chars=30000)}",
        ]
    )


def latex_escape(text: str) -> str:
    """转义 LaTeX 特殊字符，避免标题、摘要或正文导致编译失败。"""
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    return "".join(replacements.get(char, char) for char in str(text))


def clean_export_body_text(text: str, max_chars: int = 24000) -> str:
    """生成正式报告正文预览，去掉解析页码和过多技术痕迹。"""
    text = format_extracted_text_for_reading(text, max_chars=max_chars)
    lines = []
    for line in text.splitlines():
        cleaned = line.strip()
        if not cleaned:
            lines.append("")
            continue
        if re.fullmatch(r"(第\s*\d+\s*页|Page\s*\d+|Slide\s*\d+)\s*[:：]?", cleaned, re.IGNORECASE):
            continue
        cleaned = re.sub(r"^(第\s*\d+\s*页|Page\s*\d+|Slide\s*\d+)\s*[:：]?\s*", "", cleaned, flags=re.IGNORECASE)
        if cleaned:
            lines.append(cleaned)
    return clean_text("\n".join(lines))


def latex_paragraphs(text: str) -> str:
    """把普通文本转成 LaTeX 段落。"""
    paragraphs = [latex_escape(item.strip()) for item in re.split(r"\n\s*\n", clean_text(text)) if item.strip()]
    return "\n\n\\par\n".join(paragraphs) or "暂无可展示正文。"


def render_latex_template(context: dict[str, str]) -> str:
    """用简单占位符渲染项目内置 LaTeX 模板。"""
    template = REPORT_TEMPLATE_PATH.read_text(encoding="utf-8")
    for key, value in context.items():
        template = template.replace("{{ " + key + " }}", value)
    return template


def build_latex_frequency_table(freq_df: pd.DataFrame, top_n: int = 15) -> str:
    """生成适合 LaTeX 报告的高频词表。"""
    if not isinstance(freq_df, pd.DataFrame) or freq_df.empty:
        return "暂无可展示词频数据。"
    word_col = "词语" if "词语" in freq_df.columns else "word"
    count_col = "频次" if "频次" in freq_df.columns else "count"
    rows = []
    for _, row in freq_df.head(top_n).iterrows():
        rows.append(f"{latex_escape(row[word_col])} & {latex_escape(row[count_col])}\\\\")
    return "\n".join([
        r"\begin{center}",
        r"\begin{tabular}{p{0.5\linewidth}r}",
        r"\toprule",
        r"词语 & 频次\\",
        r"\midrule",
        *rows,
        r"\bottomrule",
        r"\end{tabular}",
        r"\end{center}",
    ])


def save_report_images(result: dict, temp_path: Path) -> list[tuple[str, Path]]:
    """生成报告中使用的系统分析图，不尝试复刻原文复杂图表。"""
    freq_df = result.get("freq_df")
    if not isinstance(freq_df, pd.DataFrame) or freq_df.empty:
        return []

    images: list[tuple[str, Path]] = []
    word_col = "词语" if "词语" in freq_df.columns else "word"
    count_col = "频次" if "频次" in freq_df.columns else "count"
    chart_df = freq_df.sort_values(count_col, ascending=False).head(12)
    font_prop = font_manager.FontProperties(fname=resolve_cjk_font_path()) if resolve_cjk_font_path() else None

    fig, ax = plt.subplots(figsize=(8.2, 4.2), facecolor="white")
    ax.barh(chart_df[word_col].astype(str)[::-1], chart_df[count_col][::-1], color="#2563eb")
    ax.set_xlabel("频次", fontproperties=font_prop)
    ax.set_title("高频词柱状图", fontproperties=font_prop)
    if font_prop:
        for label in ax.get_yticklabels() + ax.get_xticklabels():
            label.set_fontproperties(font_prop)
    ax.grid(axis="x", alpha=0.18)
    fig.tight_layout()
    bar_path = temp_path / "freq_bar.png"
    fig.savefig(bar_path, dpi=180, bbox_inches="tight")
    plt.close(fig)
    images.append(("高频词柱状图", bar_path))

    cloud_path = temp_path / "wordcloud.png"
    cloud_bytes = figure_to_png_bytes(build_wordcloud(freq_df, top_n=45))
    cloud_path.write_bytes(cloud_bytes)
    images.append(("高频词词云图", cloud_path))
    return images


def build_latex_chart_blocks(images: list[tuple[str, Path]]) -> str:
    """生成 LaTeX 图片块。"""
    blocks = []
    for title, image_path in images:
        blocks.append(
            "\n".join([
                r"\begin{figure}[H]",
                r"\centering",
                rf"\includegraphics[width=0.92\linewidth]{{{image_path.name}}}",
                rf"\caption{{{latex_escape(title)}}}",
                r"\end{figure}",
            ])
        )
    return "\n\n".join(blocks)


def build_latex_report_source(result: dict, temp_path: Path | None = None) -> str:
    """按项目内置 LaTeX 模板生成报告源码。"""
    chinese_summary, english_summary = split_bilingual_summary(str(result["summary"]))
    cn_chars, en_words = count_chinese_and_english(str(result["summary"]))
    text_cn_chars, text_en_words = count_chinese_and_english(str(result["extracted_text"]))
    freq_df = result.get("freq_df")
    keywords = ""
    if isinstance(freq_df, pd.DataFrame) and not freq_df.empty:
        word_col = "词语" if "词语" in freq_df.columns else "word"
        keywords = "；".join(str(row[word_col]) for _, row in freq_df.head(8).iterrows())
    english_summary_block = ""
    if english_summary:
        english_summary_block = (
            r"\begin{abstract}" + "\n"
            + r"\noindent\textbf{English Summary.} "
            + latex_escape(english_summary) + "\n"
            + r"\end{abstract}"
        )
    images = save_report_images(result, temp_path) if temp_path is not None else []
    return render_latex_template(
        {
            "filename": latex_escape(result["filename"]),
            "category": latex_escape(result["category"]),
            "chinese_summary": latex_escape(chinese_summary),
            "english_summary_block": english_summary_block,
            "keywords": latex_escape(keywords or str(result["category"])),
            "summary_stats": latex_escape(f"中文 {cn_chars} 字 / 英文 {en_words} 词"),
            "text_stats": latex_escape(f"中文 {text_cn_chars} 字 / 英文 {text_en_words} 词"),
            "category_explanation": latex_paragraphs(
                build_category_explanation(
                    str(result["category"]),
                    result.get("category_scores", {}),
                    str(result["extracted_text"]),
                )
            ),
            "frequency_table": build_latex_frequency_table(freq_df),
            "chart_blocks": build_latex_chart_blocks(images),
            "body_text": latex_paragraphs(clean_export_body_text(str(result["extracted_text"]), max_chars=24000)),
        }
    )


def build_latex_pdf_export(result: dict) -> bytes | None:
    """优先用 XeLaTeX 编译论文式 PDF；环境缺失时返回 None 走兜底方案。"""
    xelatex = shutil.which("xelatex")
    if not xelatex:
        return None
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        tex_path = temp_path / "report.tex"
        pdf_path = temp_path / "report.pdf"
        tex_path.write_text(build_latex_report_source(result, temp_path=temp_path), encoding="utf-8")
        for _ in range(2):
            completed = subprocess.run(
                [xelatex, "-interaction=nonstopmode", "-halt-on-error", str(tex_path.name)],
                cwd=temp_path,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=60,
                check=False,
            )
            if completed.returncode != 0:
                return None
        return pdf_path.read_bytes() if pdf_path.exists() else None


def set_docx_run_font(run, font_name: str = "宋体", size: float = 10.5, bold: bool = False) -> None:
    """设置 Word run 的中西文字体，避免中文回退成默认字体。"""
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    run.font.size = Pt(size)
    run.bold = bold


def set_docx_paragraph_font(paragraph, font_name: str = "宋体", size: float = 10.5, bold: bool = False) -> None:
    """批量设置段落字体。"""
    for run in paragraph.runs:
        set_docx_run_font(run, font_name=font_name, size=size, bold=bold)


def build_docx_export(result: dict) -> bytes:
    """导出排版更完整的 Word 分析报告。"""
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(1.0)
    section.left_margin = Inches(1.0)
    section.right_margin = Inches(1.0)

    styles = document.styles
    styles["Normal"].font.name = "宋体"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    styles["Normal"].font.size = Pt(10.5)

    title = document.add_heading("文档智能分析报告", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_docx_paragraph_font(title, font_name="黑体", size=18, bold=True)
    subtitle = document.add_paragraph(str(result["filename"]))
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_docx_paragraph_font(subtitle, font_name="宋体", size=10.5)

    cn_chars, en_words = count_chinese_and_english(str(result["summary"]))
    text_cn_chars, text_en_words = count_chinese_and_english(str(result["extracted_text"]))
    info_table = document.add_table(rows=5, cols=2)
    info_table.style = "Table Grid"
    info_items = [
        ("文档分类", str(result["category"])),
        ("自动分类依据", build_category_explanation(str(result["category"]), result.get("category_scores", {}), str(result["extracted_text"]))),
        ("中文摘要字数 / 英文摘要词数", f"{cn_chars} 字 / {en_words} 词"),
        ("原文中文字符 / 英文单词", f"{text_cn_chars} 字 / {text_en_words} 词"),
        ("导出说明", "摘要、分类、词频与原文均来自系统自动解析结果，建议正式提交前人工复核。"),
    ]
    for row, (label, value) in zip(info_table.rows, info_items):
        row.cells[0].text = label
        row.cells[1].text = value
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                set_docx_paragraph_font(paragraph, font_name="宋体", size=10)

    heading = document.add_heading("摘要", level=2)
    set_docx_paragraph_font(heading, font_name="黑体", size=14, bold=True)
    chinese_summary, english_summary = split_bilingual_summary(str(result["summary"]))
    p = document.add_paragraph(chinese_summary)
    p.paragraph_format.first_line_indent = Inches(0.32)
    p.paragraph_format.line_spacing = 1.5
    set_docx_paragraph_font(p, font_name="宋体", size=10.5)
    if english_summary:
        heading = document.add_heading("English Summary", level=2)
        set_docx_paragraph_font(heading, font_name="Times New Roman", size=14, bold=True)
        p = document.add_paragraph(english_summary)
        p.paragraph_format.line_spacing = 1.5
        set_docx_paragraph_font(p, font_name="Times New Roman", size=10.5)

    freq_df = result.get("freq_df")
    if isinstance(freq_df, pd.DataFrame) and not freq_df.empty:
        heading = document.add_heading("高频词", level=2)
        set_docx_paragraph_font(heading, font_name="黑体", size=14, bold=True)
        word_col = "词语" if "词语" in freq_df.columns else "word"
        count_col = "频次" if "频次" in freq_df.columns else "count"
        freq_table = document.add_table(rows=1, cols=2)
        freq_table.style = "Table Grid"
        freq_table.rows[0].cells[0].text = "词语"
        freq_table.rows[0].cells[1].text = "频次"
        for _, row in freq_df.head(15).iterrows():
            cells = freq_table.add_row().cells
            cells[0].text = str(row[word_col])
            cells[1].text = str(row[count_col])
        for table_row in freq_table.rows:
            for cell in table_row.cells:
                for paragraph in cell.paragraphs:
                    set_docx_paragraph_font(paragraph, font_name="宋体", size=10)
        with TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            for title_text, image_path in save_report_images(result, temp_path):
                p = document.add_paragraph(title_text)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                set_docx_paragraph_font(p, font_name="宋体", size=10, bold=True)
                document.add_picture(str(image_path), width=Inches(5.9))

    heading = document.add_heading("正文预览", level=2)
    set_docx_paragraph_font(heading, font_name="黑体", size=14, bold=True)
    formatted_text = clean_export_body_text(str(result["extracted_text"]), max_chars=24000)
    for paragraph in formatted_text.splitlines():
        if paragraph.strip():
            p = document.add_paragraph(paragraph.strip())
            p.paragraph_format.first_line_indent = Inches(0.32)
            p.paragraph_format.line_spacing = 1.5
            set_docx_paragraph_font(p, font_name="宋体", size=10.5)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pdf_export(result: dict) -> bytes:
    """导出论文式 PDF；优先使用 LaTeX，失败时使用 PyMuPDF 兜底。"""
    latex_pdf = build_latex_pdf_export(result)
    if latex_pdf:
        return latex_pdf

    doc = fitz.open()
    resolved_font = resolve_cjk_font_path()
    font_path = Path(resolved_font) if resolved_font else None
    font_name = "helv"
    if font_path:
        font_name = "custom"

    def wrap_text(text: str, limit: int = 42) -> list[str]:
        wrapped = []
        for raw_line in str(text).splitlines():
            line = raw_line.strip()
            if not line:
                wrapped.append("")
                continue
            while len(line) > limit:
                wrapped.append(line[:limit])
                line = line[limit:]
            wrapped.append(line)
        return wrapped

    cn_chars, en_words = count_chinese_and_english(str(result["summary"]))
    text_cn_chars, text_en_words = count_chinese_and_english(str(result["extracted_text"]))
    chinese_summary, english_summary = split_bilingual_summary(str(result["summary"]))
    sections = [
        ("基础信息", [
            f"文件名：{result['filename']}",
            f"文档分类：{result['category']}",
            f"摘要统计：中文 {cn_chars} 字 / 英文 {en_words} 词",
            f"原文统计：中文 {text_cn_chars} 字 / 英文 {text_en_words} 词",
        ]),
        ("中文摘要", wrap_text(chinese_summary)),
    ]
    if english_summary:
        sections.append(("English Summary", wrap_text(english_summary, limit=68)))
    freq_df = result.get("freq_df")
    if isinstance(freq_df, pd.DataFrame) and not freq_df.empty:
        word_col = "词语" if "词语" in freq_df.columns else "word"
        count_col = "频次" if "频次" in freq_df.columns else "count"
        freq_lines = [f"{row[word_col]}：{row[count_col]}" for _, row in freq_df.head(15).iterrows()]
        sections.append(("词频分析", freq_lines))
    sections.append(("正文预览", wrap_text(clean_export_body_text(str(result["extracted_text"]), max_chars=18000))))

    page = doc.new_page(width=595, height=842)
    if font_path:
        page.insert_font(fontname=font_name, fontfile=str(font_path))
    x, y = 50, 50
    page.insert_text((x, y), "文档智能分析报告", fontsize=18, fontname=font_name)
    y += 34
    for heading, body_lines in sections:
        if y > 760:
            page = doc.new_page(width=595, height=842)
            if font_path:
                page.insert_font(fontname=font_name, fontfile=str(font_path))
            y = 50
        page.insert_text((x, y), heading, fontsize=13, fontname=font_name)
        y += 22
        for line in body_lines:
            if y > 790:
                page = doc.new_page(width=595, height=842)
                if font_path:
                    page.insert_font(fontname=font_name, fontfile=str(font_path))
                y = 50
            page.insert_text((x, y), line, fontsize=10.5, fontname=font_name)
            y += 16
        y += 10
    data = doc.tobytes()
    doc.close()
    return data


def build_download_payload(result: dict, export_format: str) -> tuple[bytes, str, str]:
    """根据用户选择构建下载内容、文件名和 MIME。"""
    stem = safe_download_stem(result["filename"])
    freq_df = result["freq_df"]
    if freq_df.empty and export_format in {"词频CSV", "词云PNG", "词频柱状图HTML"}:
        return "暂无可下载的词频数据。".encode("utf-8"), f"{stem}_词频为空.txt", "text/plain"
    score_df = pd.DataFrame(
        sorted(result["category_scores"].items(), key=lambda item: item[1], reverse=True),
        columns=["类别", "得分"],
    )
    if export_format == "摘要TXT":
        chinese_summary, english_summary = split_bilingual_summary(str(result["summary"]))
        content = [f"中文摘要：\n{chinese_summary}"]
        if english_summary:
            content.append(f"\nEnglish Summary:\n{english_summary}")
        return "\n".join(content).encode("utf-8"), f"{stem}_摘要.txt", "text/plain"
    if export_format == "词频CSV":
        return freq_df.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"), f"{stem}_词频.csv", "text/csv"
    if export_format == "分类得分CSV":
        return score_df.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig"), f"{stem}_分类得分.csv", "text/csv"
    if export_format == "词云PNG":
        return figure_to_png_bytes(build_wordcloud(freq_df, top_n=int(st.session_state.get("word_top_n", 30)))), f"{stem}_词云.png", "image/png"
    if export_format == "词频柱状图HTML":
        html = build_word_frequency_chart(freq_df, top_n=int(st.session_state.get("word_top_n", 30))).to_html(include_plotlyjs="cdn", full_html=True)
        return html.encode("utf-8"), f"{stem}_词频柱状图.html", "text/html"
    if export_format == "Markdown":
        return build_markdown_export(result).encode("utf-8"), f"{stem}_分析结果.md", "text/markdown"
    if export_format == "Word":
        return build_docx_export(result), f"{stem}_分析结果.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if export_format == "PDF":
        return build_pdf_export(result), f"{stem}_分析结果.pdf", "application/pdf"
    return result["extracted_text"].encode("utf-8"), f"{stem}_提取文本.txt", "text/plain"


def reset_upload_state(clear_result: bool = True) -> None:
    """重置上传控件和本次处理状态，用于清空错误文件或重新选择文件。"""
    st.session_state.upload_key = f"upload_{uuid.uuid4().hex}"
    st.session_state.processed_upload = None
    st.session_state.upload_error = None
    if clear_result:
        st.session_state.last_result = None


# ---------- Streamlit 页面入口 ----------

st.set_page_config(page_title="文档智能处理系统", page_icon="📄", layout="wide")

with st.spinner("加载中，请稍候..."):
    init_storage()
    # 初始化 Streamlit 会话状态，避免页面刷新后丢失登录用户和最近一次分析结果。
    if "user" not in st.session_state:
        st.session_state.user = None
    if "last_result" not in st.session_state:
        st.session_state.last_result = None
    if "upload_key" not in st.session_state:
        st.session_state.upload_key = f"upload_{uuid.uuid4().hex}"
    if "processed_upload" not in st.session_state:
        st.session_state.processed_upload = None
    if "upload_error" not in st.session_state:
        st.session_state.upload_error = None

st.title("文档智能处理系统")
st.caption("一个支持多用户、文档解析、OCR、摘要生成、词频可视化和自动分类的在线文档处理小工具。")

st.markdown(
    """
    <style>
    .block-container {
        padding-top: 2.6rem;
        padding-bottom: 2rem;
        max-width: 1280px;
    }
    h1 {
        line-height: 1.25 !important;
        margin-top: 0.25rem !important;
        margin-bottom: 0.6rem !important;
        overflow: visible !important;
    }
    .stApp {
        background: #f6f7fb;
        color: #111827;
        color-scheme: light;
    }
    html, body, [class*="css"] {
        color: #111827 !important;
        color-scheme: light;
    }
    h1, h2, h3, h4, h5, h6, p, label, span, div {
        color: #111827;
    }
    [data-testid="stSidebar"] {
        background: #ffffff !important;
        color: #111827 !important;
    }
    [data-testid="stHeader"] {
        background: #f6f7fb !important;
    }
    [data-testid="stHeader"] *,
    [data-testid="stToolbar"] *,
    [data-testid="stDecoration"] *,
    [data-testid="stStatusWidget"] *,
    [data-testid="stMainMenu"] *,
    [data-testid="stDeployButton"] *,
    [data-testid="baseButton-header"] *,
    [data-testid="collapsedControl"] * {
        color: #111827 !important;
        fill: #111827 !important;
        stroke: #111827 !important;
    }
    [data-testid="stToolbar"] button,
    [data-testid="stStatusWidget"] button,
    [data-testid="stMainMenu"] button,
    [data-testid="stDeployButton"] button,
    [data-testid="baseButton-header"] {
        background: rgba(255, 255, 255, 0.92) !important;
        color: #111827 !important;
        border-color: #d1d5db !important;
    }
    [data-testid="stToolbar"] svg,
    [data-testid="stStatusWidget"] svg,
    [data-testid="stMainMenu"] svg,
    [data-testid="stDeployButton"] svg,
    [data-testid="baseButton-header"] svg {
        color: #111827 !important;
        fill: #111827 !important;
        stroke: #111827 !important;
    }
    div[data-testid="stAlert"], div[data-testid="stExpander"], div[data-testid="stDataFrame"],
    div[data-testid="stForm"], div[data-testid="stFileUploader"], div[data-testid="stTabs"],
    div[data-testid="stMarkdownContainer"] {
        color: #111827 !important;
    }
    textarea, input, select, .stTextInput input, .stTextArea textarea {
        background: #ffffff !important;
        color: #111827 !important;
        border-color: #d1d5db !important;
    }
    button, button * {
        color: #111827 !important;
    }
    .stButton button, .stDownloadButton button, div[data-testid="stFormSubmitButton"] button {
        background: #ffffff !important;
        color: #111827 !important;
        border: 1px solid #d1d5db !important;
    }
    .stButton button[kind="primary"], div[data-testid="stFormSubmitButton"] button[kind="primary"] {
        background: #2563eb !important;
        color: #ffffff !important;
        border-color: #2563eb !important;
    }
    div[data-testid="stTabs"] button, div[data-testid="stTabs"] button p {
        color: #111827 !important;
    }
    div[data-testid="stFileUploaderDropzone"] {
        background: #fbfcfe !important;
        color: #111827 !important;
        border: 1px dashed #cbd5e1 !important;
        border-radius: 8px !important;
        box-shadow: none;
        transition: background 0.15s ease, border-color 0.15s ease, box-shadow 0.15s ease;
    }
    div[data-testid="stFileUploaderDropzone"]:hover {
        background: #f8fafc !important;
        border-color: #94a3b8 !important;
        box-shadow: 0 1px 3px rgba(15, 23, 42, 0.06);
    }
    div[data-testid="stFileUploaderDropzone"] * {
        color: #111827 !important;
    }
    div[data-testid="stFileUploaderDropzone"] small,
    div[data-testid="stFileUploaderDropzone"] [data-testid="stMarkdownContainer"] p {
        display: none !important;
    }
    div[data-testid="stFileUploader"] section {
        background: transparent !important;
    }
    .upload-hint {
        background: #ffffff;
        border: 1px solid #e5e7eb;
        border-radius: 8px;
        padding: 0.55rem 0.7rem;
        margin: 0.3rem 0 0.45rem 0;
        color: #4b5563;
        font-size: 0.88rem;
        line-height: 1.35;
    }
    [data-baseweb="select"] *, [data-baseweb="radio"] *, [data-baseweb="slider"] * {
        color: #111827 !important;
    }
    div[data-testid="stMetric"] {
        background: #ffffff;
        border: 1px solid #e5e7eb;
        border-radius: 10px;
        padding: 0.45rem 0.7rem;
    }
    div[data-testid="stMetric"] label {
        color: #6b7280;
        font-size: 0.82rem;
    }
    div[data-testid="stMetric"] [data-testid="stMetricValue"] {
        font-size: 1.1rem;
        line-height: 1.2;
    }
    .feature-strip {
        display: flex;
        gap: 0.75rem;
        margin: 0.25rem 0 0.5rem 0;
        flex-wrap: wrap;
    }
    .feature-chip {
        flex: 1 1 180px;
        background: #ffffff;
        border: 1px solid #e5e7eb;
        border-radius: 10px;
        padding: 0.65rem 0.8rem;
        box-shadow: 0 1px 2px rgba(0, 0, 0, 0.03);
    }
    .feature-chip .label {
        color: #6b7280;
        font-size: 0.78rem;
        line-height: 1.1;
        margin-bottom: 0.25rem;
    }
    .feature-chip .value {
        color: #111827;
        font-size: 0.98rem;
        font-weight: 600;
        line-height: 1.2;
    }
    .entry-note {
        background: #ffffff;
        border: 1px solid #e5e7eb;
        border-radius: 8px;
        padding: 0.85rem 1rem;
        margin: 0.55rem 0 1rem 0;
        color: #374151;
        line-height: 1.65;
    }
    .entry-note strong {
        color: #111827;
    }
    div[data-testid="stTabs"] button {
        border-radius: 8px;
    }
    .stSelectbox label, .stTextInput label, .stRadio label, .stFileUploader label {
        font-weight: 600;
    }
    @media (max-width: 760px) {
        .block-container {
            padding-top: 2.1rem;
            padding-left: 0.85rem;
            padding-right: 0.85rem;
        }
        h1 {
            font-size: 1.65rem !important;
            line-height: 1.25 !important;
        }
        .feature-strip {
            gap: 0.45rem;
        }
        .feature-chip {
            flex: 1 1 100%;
            padding: 0.5rem 0.65rem;
        }
        .feature-chip .value {
            font-size: 0.9rem;
        }
    }
    </style>
    """,
    unsafe_allow_html=True,
)

st.markdown(
    """
    <div class="feature-strip">
      <div class="feature-chip">
        <div class="label">支持</div>
        <div class="value">PDF / Word / 扫描件OCR</div>
      </div>
      <div class="feature-chip">
        <div class="label">输出</div>
        <div class="value">摘要 / 分类 / 词云词频</div>
      </div>
      <div class="feature-chip">
        <div class="label">历史</div>
        <div class="value">当前账号隔离保存</div>
      </div>
      <div class="feature-chip">
        <div class="label">部署</div>
        <div class="value">公网链接可访问</div>
      </div>
    </div>
    """,
    unsafe_allow_html=True,
)

# 未登录时只显示登录和注册表单；登录成功后再进入主功能页面。
if st.session_state.user is None:
    st.markdown(
        """
        <div class="entry-note">
            <strong>评阅入口：</strong>可直接点击“免注册试用”进入系统，上传示例文件检查解析、OCR、摘要、词频图和分类结果。
            也可以注册新账号；不同账号的上传记录互相隔离，只能看到自己的文件。
        </div>
        """,
        unsafe_allow_html=True,
    )
    if st.button("免注册试用，直接进入系统", width="stretch"):
        delete_demo_documents()
        st.session_state.user = get_or_create_demo_user()
        st.session_state.last_result = None
        st.rerun()

    login_tab, register_tab = st.tabs(["已有账号登录", "新用户注册"])

    with login_tab:
        with st.form("login_form"):
            login_username = st.text_input("用户名", key="login_username")
            login_password = st.text_input("密码", type="password", key="login_password")
            submitted = st.form_submit_button("登录")
        if submitted:
            user = authenticate(login_username, login_password)
            if user is None:
                st.error("用户名或密码错误。")
            else:
                st.session_state.user = user
                st.session_state.last_result = None
                st.rerun()

    with register_tab:
        with st.form("register_form"):
            register_username = st.text_input("用户名", key="register_username")
            register_password = st.text_input("密码", type="password", key="register_password")
            register_password_2 = st.text_input("确认密码", type="password", key="register_password_2")
            submitted = st.form_submit_button("注册")
        if submitted:
            if register_password != register_password_2:
                st.error("两次输入的密码不一致。")
            else:
                ok, message = create_user(register_username, register_password)
                if ok:
                    user = authenticate(register_username, register_password)
                    if user is not None:
                        st.session_state.user = user
                        st.session_state.last_result = None
                        st.rerun()
                    st.success(message)
                else:
                    st.error(message)
    st.stop()

user = st.session_state.user
if is_demo_user(user):
    delete_demo_documents()

# 侧边栏负责展示 OCR 状态、退出登录按钮和当前用户的历史文档。
with st.sidebar:
    st.write(f"当前用户：{user['username']}")
    ocr_ready, ocr_message = get_tesseract_status()
    if ocr_ready:
        st.caption(f"OCR 已就绪：{ocr_message}")
    else:
        st.warning(ocr_message)
    if st.button("退出登录"):
        st.session_state.user = None
        st.session_state.last_result = None
        st.rerun()

    with st.expander("处理模式", expanded=True):
        preset_label = st.radio(
            "我想要",
            ["日常快速", "论文精修", "扫描件识别"],
            index=0,
            help="普通用户优先选这里即可；只有需要更细控制时再展开手动微调。",
        )
        preset_defaults = {
            "日常快速": {
                "ocr_mode": "standard",
                "pdf_extract_mode": "auto",
                "summary_mode": "中文摘要",
                "quality_mode": "快速",
                "pdca_cycles": 3,
                "pdca_pass_score": 88,
                "cn_target": 320,
                "en_target": 180,
            },
            "论文精修": {
                "ocr_mode": "standard",
                "pdf_extract_mode": "auto",
                "summary_mode": "双语摘要",
                "quality_mode": "精修",
                "pdca_cycles": 3,
                "pdca_pass_score": 90,
                "cn_target": 360,
                "en_target": 220,
            },
            "扫描件识别": {
                "ocr_mode": "accurate",
                "pdf_extract_mode": "vision",
                "summary_mode": "中文摘要",
                "quality_mode": "快速",
                "pdca_cycles": 3,
                "pdca_pass_score": 85,
                "cn_target": 280,
                "en_target": 160,
            },
        }
        preset = preset_defaults[preset_label]
        preset_tips = {
            "日常快速": "适合普通 Word/PDF，默认三轮复核，兼顾速度和稳妥。",
            "论文精修": "适合论文、项目材料和给老师看的摘要，会更慢但摘要更正式。",
            "扫描件识别": "适合图片或扫描版 PDF，会加强 OCR 识别并默认三轮复核，处理时间更长。",
        }
        st.caption(preset_tips[preset_label])

        with st.expander("手动微调", expanded=False):
            ocr_options = ["快速", "标准", "高精度"]
            ocr_default = {"fast": 0, "standard": 1, "accurate": 2}[preset["ocr_mode"]]
            ocr_mode_label = st.radio(
                "图片/扫描 PDF 识别",
                ocr_options,
                index=ocr_default,
                horizontal=True,
                help="普通文字 PDF 不受影响。扫描件不清楚时选高精度。",
            )
            preset["ocr_mode"] = {"快速": "fast", "标准": "standard", "高精度": "accurate"}[ocr_mode_label]

            pdf_extract_options = ["自动推荐", "快速文本", "多模态精读"]
            pdf_extract_default = {"auto": 0, "text": 1, "vision": 2}[preset["pdf_extract_mode"]]
            pdf_extract_label = st.radio(
                "PDF提取方式",
                pdf_extract_options,
                index=pdf_extract_default,
                horizontal=True,
                help="自动推荐会先用快速文本提取，质量差时再尝试多模态精读。",
            )
            preset["pdf_extract_mode"] = {"自动推荐": "auto", "快速文本": "text", "多模态精读": "vision"}[pdf_extract_label]

            summary_options = ["中文摘要", "英文摘要", "双语摘要"]
            summary_mode = st.radio(
                "摘要语言",
                summary_options,
                index=summary_options.index(preset["summary_mode"]),
                horizontal=True,
                help="给中文材料通常选中文摘要；论文展示可选双语摘要。",
            )
            preset["summary_mode"] = summary_mode

            quality_mode = st.radio(
                "摘要质量",
                ["快速", "精修"],
                index=0 if preset["quality_mode"] == "快速" else 1,
                horizontal=True,
                help="精修会额外复核和润色，质量更高但更慢。",
            )
            preset["quality_mode"] = quality_mode

            preset["pdca_cycles"] = st.slider(
                "质量复核轮次",
                min_value=1,
                max_value=5,
                value=int(preset["pdca_cycles"]),
                help="轮次越多越慢；日常使用 1 轮即可。",
            )
            preset["pdca_pass_score"] = st.slider(
                "复核通过分数",
                min_value=70,
                max_value=95,
                value=int(preset["pdca_pass_score"]),
                step=1,
                help="越高越严格，也越容易触发修订。",
            )
            preset["cn_target"] = st.slider(
                "中文摘要字数",
                min_value=120,
                max_value=800,
                value=int(preset["cn_target"]),
                step=20,
            )
            preset["en_target"] = st.slider(
                "英文摘要词数",
                min_value=80,
                max_value=500,
                value=int(preset["en_target"]),
                step=10,
            )
        st.session_state.ocr_mode = preset["ocr_mode"]
        st.session_state.pdf_extract_mode = preset["pdf_extract_mode"]
        st.session_state.summary_mode = preset["summary_mode"]
        st.session_state.quality_mode = preset["quality_mode"]
        st.session_state.pdca_cycles = int(preset["pdca_cycles"])
        st.session_state.pdca_pass_score = int(preset["pdca_pass_score"])
        st.session_state.cn_summary_target = int(preset["cn_target"])
        st.session_state.en_summary_target = int(preset["en_target"])

        category_mode = st.radio(
            "文档分类方式",
            ["自动分类", "选择标准分类", "自定义分类"],
            index=0,
            help="自动分类会根据文档内容判断；手动分类会保存到历史记录，便于后续筛选。",
        )
        st.session_state.category_mode = category_mode
        if category_mode == "选择标准分类":
            st.session_state.manual_category = st.selectbox("标准分类", STANDARD_CATEGORIES)
            st.session_state.custom_category = ""
        elif category_mode == "自定义分类":
            st.session_state.custom_category = st.text_input("新建分类名称", placeholder="例如：课程材料、企业案例、导师资料")
            st.session_state.manual_category = ""
        else:
            st.session_state.manual_category = ""
            st.session_state.custom_category = ""
        word_top_n = st.slider("词频显示数量", min_value=10, max_value=60, value=30, step=5)
        st.session_state.word_top_n = word_top_n
        st.caption(
            f"当前：{preset_label}｜"
            f"{'双语' if preset['summary_mode'] == '双语摘要' else preset['summary_mode']}｜"
            f"{preset['quality_mode']}｜复核{preset['pdca_cycles']}轮"
        )

    st.divider()
    st.subheader("我的文档")
    documents = list_user_documents(user["id"])
    if not documents:
        st.caption("当前账号暂无上传记录。")
        selected_doc_id = None
    else:
        saved_categories = {doc["category"] for doc in documents if doc["category"]}
        standard_saved = [category for category in STANDARD_CATEGORIES if category in saved_categories]
        custom_saved = sorted(saved_categories - set(STANDARD_CATEGORIES))
        categories = ["全部分类", *standard_saved, *custom_saved]
        selected_category = st.radio("按分类筛选", categories, horizontal=False, index=0)
        filtered_documents = [doc for doc in documents if selected_category == "全部分类" or doc["category"] == selected_category]
        st.caption(f"共 {len(filtered_documents)} 条记录")
        if not filtered_documents:
            st.caption("当前分类下暂无文档。")
            selected_doc_id = None
        else:
            options = {
                f"{doc['created_at']}  ·  {format_duration(doc['processing_seconds'])}  ·  {doc['filename']}": doc["id"]
                for doc in filtered_documents
            }
            selected_label = st.selectbox("最近文档", list(options.keys()))
            selected_doc_id = options[selected_label]
        if selected_doc_id is not None and st.button("查看历史文档"):
            doc = get_user_document(user["id"], selected_doc_id)
            if doc is not None:
                auto_category, category_scores = classify_document(doc["extracted_text"])
                route = route_file(doc["filename"], doc["extracted_text"]).to_dict()
                category = doc["category"] or route.get("task_label") or auto_category
                freq = Counter(tokenize(doc["extracted_text"])).most_common(30)
                freq_df = pd.DataFrame(freq, columns=["word", "count"])
                st.session_state.last_result = {
                    "filename": doc["filename"],
                    "category": category,
                    "category_scores": category_scores,
                    "summary": doc["summary"],
                    "summary_source": "历史保存摘要",
                    "pdca_report": doc["pdca_report"] or "该历史记录未保存PDCA检查报告。",
                    "freq_df": freq_df,
                    "extracted_text": doc["extracted_text"],
                    "processing_seconds": float(doc["processing_seconds"] or 0),
                    "route": route,
                    "from_history": True,
                }
                st.rerun()

uploaded_file = st.file_uploader(
    "上传文档并开始处理",
    type=["pdf", "doc", "docx", "ppt", "pptx", "csv", "xlsx", "xls", "txt", "md", "json", "jsonl", "png", "jpg", "jpeg", "bmp", "tif", "tiff"],
    key=st.session_state.upload_key,
)
st.markdown(
    '<div class="upload-hint">单个文件最大 200MB；支持 PDF、Word、扫描版 PDF、图片、PPT、表格和文本。选择文件后会自动提取文本、生成摘要、分类并统计词频。</div>',
    unsafe_allow_html=True,
)

if st.button("重新选择文件"):
    reset_upload_state(clear_result=False)
    st.rerun()
st.caption("历史记录会自动按当前账号保存。")

if st.session_state.upload_error:
    st.error(st.session_state.upload_error)
    st.info("请点击“清空本次上传”后重新选择正确格式的文件。")
    st.stop()

if uploaded_file is None:
    st.session_state.upload_error = None

if uploaded_file is not None:
    suffix = uploaded_file.name.lower().rsplit(".", 1)[-1] if "." in uploaded_file.name else ""
    allowed_suffixes = {"pdf", "doc", "docx", "ppt", "pptx", "csv", "xlsx", "xls", "txt", "md", "json", "jsonl", "png", "jpg", "jpeg", "bmp", "tif", "tiff"}
    if suffix not in allowed_suffixes:
        st.session_state.upload_error = f"文件类型不支持：{uploaded_file.name}。请上传 PDF、Word、PPT、表格、文本或图片文件。"
        st.session_state.processed_upload = None
        st.rerun()

if uploaded_file is not None:
    uploaded_file_hash = hashlib.sha256(uploaded_file.getvalue()).hexdigest()
    upload_signature = (
        f"{st.session_state.upload_key}:{uploaded_file.name}:{uploaded_file.size}:{uploaded_file_hash}:"
        f"{st.session_state.get('ocr_mode', 'standard')}:"
        f"{st.session_state.get('pdf_extract_mode', 'auto')}:"
        f"{st.session_state.get('summary_mode', '中文摘要')}:"
        f"{st.session_state.get('quality_mode', '快速')}:"
        f"{st.session_state.get('pdca_cycles', 3)}:"
        f"{st.session_state.get('pdca_pass_score', 88)}:"
        f"{st.session_state.get('category_mode', '自动分类')}:"
        f"{st.session_state.get('manual_category', '')}:"
        f"{st.session_state.get('custom_category', '')}"
    )
else:
    upload_signature = None

# 只有当前文件尚未处理过时才执行解析，避免 Streamlit 重跑导致重复保存。
if uploaded_file is not None and st.session_state.processed_upload != upload_signature:
    processing_started_at = time.perf_counter()
    file_bytes = uploaded_file.getvalue()
    st.session_state.upload_error = None
    progress_bar = st.progress(0)
    status_box = st.empty()
    eta_hint = estimate_processing_time(
        uploaded_file.name,
        uploaded_file.size,
        st.session_state.get("ocr_mode", "standard"),
        st.session_state.get("summary_mode", "中文摘要"),
    )
    st.caption(f"本次处理{eta_hint}。扫描件、长 PDF、双语摘要或质量复核会更慢。")
    update_progress(progress_bar, status_box, 5, "已接收文件，准备解析", eta_hint)

    try:
        with st.spinner("正在处理文件，请稍候..."):
            # 解析阶段会根据文件类型自动选择 PDF、Word、PPT、表格文本提取或 OCR。
            update_progress(progress_bar, status_box, 18, "正在提取文本/OCR 识别", eta_hint)
            extracted_text, cache_hit = parse_uploaded_file(
                uploaded_file,
                ocr_mode=st.session_state.get("ocr_mode", "standard"),
                pdf_extract_mode=st.session_state.get("pdf_extract_mode", "auto"),
            )
            if cache_hit:
                update_progress(progress_bar, status_box, 45, "已命中解析缓存，正在检查文本质量", eta_hint)
            else:
                update_progress(progress_bar, status_box, 45, "文本提取完成，正在检查文本质量", eta_hint)
    except Exception as exc:
        progress_bar.progress(100)
        st.session_state.upload_error = (
            f"文件解析失败：{uploaded_file.name}。可能是文件扩展名和真实格式不一致、文件损坏、加密，"
            f"或该格式暂不支持。错误信息：{exc}"
        )
        st.session_state.processed_upload = upload_signature
        status_box.error(st.session_state.upload_error)
        if st.button("清空错误文件并重新上传"):
            reset_upload_state(clear_result=True)
            st.rerun()
        st.stop()

    if not extracted_text:
        progress_bar.progress(100)
        st.session_state.processed_upload = upload_signature
        st.session_state.upload_error = "未提取到有效文本。请确认文件内容清晰，扫描件建议使用高分辨率图片或 PDF。"
        status_box.warning("未提取到有效文本。")
        st.warning(st.session_state.upload_error)
        if st.button("清空无效文件并重新上传"):
            reset_upload_state(clear_result=True)
            st.rerun()
        st.stop()

    if is_unreadable_text(extracted_text):
        # 乱码文本不进入摘要流程，防止生成看似正式但实际不可靠的摘要。
        progress_bar.progress(100)
        st.session_state.processed_upload = upload_signature
        st.session_state.upload_error = "提取结果疑似存在大量乱码，已停止摘要生成。请检查文件编码、字体或重新导出为清晰 PDF/Word 后再上传。"
        status_box.warning("提取结果疑似存在大量乱码，已停止摘要生成。")
        st.warning(st.session_state.upload_error)
        st.text_area("当前提取结果", extracted_text, height=240)
        if st.button("清空乱码文件并重新上传"):
            reset_upload_state(clear_result=True)
            st.rerun()
        st.stop()

    update_progress(progress_bar, status_box, 55, "正在分析文档内容并生成摘要，请稍候", eta_hint)
    with st.spinner("正在生成分析结果，请稍候..."):
        category, category_scores, summary, summary_source, pdca_report, freq_df, route = analyze_text(
            extracted_text,
            uploaded_file.name,
            st.session_state.get("summary_mode", "中文摘要"),
            st.session_state.get("pdca_cycles", 3),
            st.session_state.get("pdca_pass_score", 88),
            st.session_state.get("cn_summary_target"),
            st.session_state.get("en_summary_target"),
            st.session_state.get("quality_mode", "快速"),
        )
        category = resolve_user_category(category)
    processing_seconds = time.perf_counter() - processing_started_at
    if is_demo_user(user):
        update_progress(progress_bar, status_box, 88, "摘要生成完成，免注册试用不保存历史记录", eta_hint)
        delete_demo_documents()
    else:
        update_progress(progress_bar, status_box, 88, "摘要生成完成，正在保存记录", eta_hint)
        # 分析结果和原始文件一起保存到当前用户目录，保证历史记录可追溯。
        save_document(
            user_id=user["id"],
            filename=uploaded_file.name,
            file_bytes=file_bytes,
            category=category,
            summary=summary,
            pdca_report=pdca_report,
            extracted_text=extracted_text,
            word_count=len(freq_df),
            processing_seconds=processing_seconds,
        )
    st.session_state.last_result = {
        "filename": uploaded_file.name,
        "category": category,
        "category_scores": category_scores,
        "summary": summary,
        "summary_source": summary_source,
        "pdca_report": pdca_report,
        "freq_df": freq_df,
        "extracted_text": extracted_text,
        "processing_seconds": processing_seconds,
        "route": route,
        "from_history": False,
    }
    st.session_state.processed_upload = upload_signature
    st.session_state.upload_error = None
    update_progress(progress_bar, status_box, 100, f"处理完成，用时 {format_duration(processing_seconds)}")
    if is_demo_user(user):
        st.success(f"文档已解析完成。本次为免注册试用，不保存系统历史记录；用时 {format_duration(processing_seconds)}。")
    else:
        st.success(f"文档已解析并保存到当前用户的历史记录。本次用时 {format_duration(processing_seconds)}。")

if st.session_state.last_result is None:
    st.info("请上传 PDF、Word、PPT、表格、文本或图片文件开始处理。上传记录只会显示在当前登录账号下。")
    st.stop()

# 页面底部展示最近一次上传或历史记录加载得到的分析结果。
result = st.session_state.last_result
summary_text = str(result["summary"])
chinese_summary, english_summary = split_bilingual_summary(summary_text)
summary_word_count = int(result["freq_df"]["count"].sum()) if not result["freq_df"].empty else 0
cn_chars, cn_words = count_chinese_and_english(chinese_summary)
en_chars, en_words = count_chinese_and_english(english_summary)
summary_top = st.columns([3, 1, 1, 1, 1])
with summary_top[0]:
    st.subheader(f"当前文档：{result['filename']}")
with summary_top[1]:
    st.metric("分类", result["category"])
with summary_top[2]:
    st.metric("总词数", summary_word_count)
with summary_top[3]:
    st.metric("中英摘要", f"{cn_chars}字 / {en_words}词")
with summary_top[4]:
    st.metric("处理用时", format_duration(result.get("processing_seconds")))

with st.expander("自动分类依据", expanded=False):
    st.write(build_category_explanation(str(result["category"]), result["category_scores"], str(result["extracted_text"])))

with st.container(border=True):
    render_summary(summary_text, result.get("summary_source", "本地摘要"))

detail_tabs = st.tabs(["分类与词频", "全文"])
with detail_tabs[0]:
    word_top_n = int(st.session_state.get("word_top_n", 30))
    left_box, right_box = st.columns([1, 2])
    with left_box:
        st.subheader("自动分类")
        st.success(result["category"])
        st.caption("分类方式：基于关键词规则打分，可按实际业务继续扩展。")
        st.metric("原文长度", f"{len(result['extracted_text'])} 字符")
        st.metric("中文摘要", f"{cn_chars} 字")
        if english_summary:
            st.metric("英文摘要", f"{en_words} 词")
    with right_box:
        st.subheader("分类得分")
        score_df = pd.DataFrame(
            sorted(result["category_scores"].items(), key=lambda item: item[1], reverse=True),
            columns=["类别", "得分"],
        )
        st.dataframe(score_df, width="stretch")
    st.subheader("词频统计")
    tab_bar, tab_cloud = st.tabs(["柱状图", "词云"])
    with tab_bar:
        fig = build_word_frequency_chart(result["freq_df"], top_n=word_top_n)
        st.plotly_chart(fig, width="stretch")
    with tab_cloud:
        try:
            st.pyplot(build_wordcloud(result["freq_df"], top_n=word_top_n), clear_figure=True)
        except Exception as exc:
            st.write(f"词云生成失败：{exc}。柱状图仍可正常使用。")
with detail_tabs[1]:
    render_extracted_text(result["extracted_text"])

with st.expander("下载与导出", expanded=False):
    st.caption("推荐优先下载 Word 或 PDF；TXT 适合复制，CSV 和图表适合二次分析。")
    primary_exports = [
        ("Word", "下载完整分析报告（Word，可编辑）"),
        ("PDF", "下载完整分析报告（论文式PDF）"),
        ("Markdown", "下载完整分析报告（Markdown文本）"),
        ("摘要TXT", "下载摘要正文（TXT）"),
    ]
    data_exports = [
        ("TXT", "下载提取全文（TXT）"),
        ("词频CSV", "下载高频词统计表（CSV）"),
        ("分类得分CSV", "下载分类得分表（CSV）"),
        ("词云PNG", "下载高频词词云图（PNG）"),
        ("词频柱状图HTML", "下载高频词柱状图（HTML）"),
    ]
    for export_row in (primary_exports, data_exports):
        columns = st.columns(len(export_row))
        for column, (export_format, label) in zip(columns, export_row):
            download_data, download_name, download_mime = build_download_payload(result, export_format)
            column.download_button(
                label,
                data=download_data,
                file_name=download_name,
                mime=download_mime,
                key=f"download_{export_format}_{safe_download_stem(result['filename'])}",
            )

